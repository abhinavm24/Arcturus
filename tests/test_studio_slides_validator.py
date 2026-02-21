"""Tests for core/studio/slides/validator.py — v2 quality checks."""

import pytest
from pathlib import Path

from pptx import Presentation

from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
from core.studio.slides.exporter import SLIDE_HEIGHT, SLIDE_WIDTH, export_to_pptx
from core.studio.slides.themes import get_theme
from core.studio.slides.validator import validate_pptx


def _export(ct, tmp_path, filename="test.pptx"):
    output = tmp_path / filename
    export_to_pptx(ct, get_theme(), output)
    return output


def _good_ct():
    return SlidesContentTree(
        deck_title="Test",
        slides=[
            Slide(id="s1", slide_type="content", title="Slide 1",
                  elements=[SlideElement(id="e1", type="body", content="Good content here with enough text.")],
                  speaker_notes="Discuss the key points of this slide. Elaborate on the main findings and their implications for the project."),
        ],
    )


def test_valid_pptx_passes(tmp_path):
    path = _export(_good_ct(), tmp_path)
    result = validate_pptx(path, expected_slide_count=1, content_tree=_good_ct())
    assert result["valid"] is True
    assert result["layout_valid"] is True


def test_invalid_file_fails(tmp_path):
    fake = tmp_path / "bad.pptx"
    fake.write_bytes(b"not a pptx")
    result = validate_pptx(fake)
    assert result["valid"] is False
    assert len(result["errors"]) > 0


def test_slide_count_mismatch(tmp_path):
    path = _export(_good_ct(), tmp_path)
    result = validate_pptx(path, expected_slide_count=10)
    assert result["valid"] is False
    assert any("mismatch" in e for e in result["errors"])


def test_block_char_overflow(tmp_path):
    ct = SlidesContentTree(
        deck_title="Overflow",
        slides=[Slide(id="s1", slide_type="content", title="X",
                      elements=[SlideElement(id="e1", type="body", content="A" * 1000)],
                      speaker_notes="Notes.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path)
    assert result["layout_valid"] is False
    assert any("800 chars" in e for e in result["layout_errors"])


def test_slide_char_overflow(tmp_path):
    ct = SlidesContentTree(
        deck_title="Dense",
        slides=[Slide(id="s1", slide_type="content", title="T" * 500,
                      elements=[SlideElement(id="e1", type="body", content="B" * 1200)],
                      speaker_notes="Notes.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path)
    assert result["layout_valid"] is False
    assert any("1600 chars" in e for e in result["layout_errors"])


def test_out_of_bounds_shape(tmp_path):
    """Shape positioned outside slide boundary detected."""
    from pptx.util import Inches
    ct = _good_ct()
    path = _export(ct, tmp_path)
    # Manually add an out-of-bounds shape
    prs = Presentation(str(path))
    slide = prs.slides[0]
    slide.shapes.add_textbox(Inches(12), Inches(6), Inches(3), Inches(3))
    prs.save(str(path))
    result = validate_pptx(path)
    assert any("beyond" in e for e in result["layout_errors"])


def test_chart_quality_valid_with_chart(tmp_path):
    """Chart slide with native chart shape passes chart quality."""
    ct = SlidesContentTree(
        deck_title="Charts",
        slides=[Slide(id="s1", slide_type="chart", title="Revenue",
                      elements=[SlideElement(id="e1", type="chart", content={
                          "chart_type": "bar",
                          "categories": ["Q1", "Q2"],
                          "series": [{"name": "Rev", "values": [1.0, 2.0]}],
                      })],
                      speaker_notes="Discuss the revenue data in detail. Highlight the key growth trend and its implications.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["chart_quality_valid"] is True


def test_chart_quality_valid_with_fallback(tmp_path):
    """Chart slide with [Chart: ...] text fallback passes chart quality."""
    ct = SlidesContentTree(
        deck_title="Charts",
        slides=[Slide(id="s1", slide_type="chart", title="Data",
                      elements=[SlideElement(id="e1", type="chart", content="Revenue growth")],
                      speaker_notes="Walk through the data trends carefully. Connect the numbers to the broader business context.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["chart_quality_valid"] is True


def test_chart_quality_invalid(tmp_path):
    """Chart slide with no chart or fallback marker fails chart quality."""
    ct = SlidesContentTree(
        deck_title="Empty Chart",
        slides=[Slide(id="s1", slide_type="chart", title="Data",
                      elements=[],
                      speaker_notes="Discuss the data. This slide should have a chart but it is missing for testing purposes.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["chart_quality_valid"] is False


def test_notes_quality_valid(tmp_path):
    good_notes = ("Discuss the key points of this slide in detail. "
                  "Highlight the main finding and explain its implications. "
                  "Transition to the next topic smoothly.")
    ct = SlidesContentTree(
        deck_title="Notes",
        slides=[Slide(id="s1", slide_type="content", title="S",
                      elements=[SlideElement(id="e1", type="body", content="Content.")],
                      speaker_notes=good_notes)],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["notes_quality_valid"] is True


def test_notes_quality_invalid(tmp_path):
    ct = SlidesContentTree(
        deck_title="No Notes",
        slides=[Slide(id="s1", slide_type="content", title="S",
                      elements=[SlideElement(id="e1", type="body", content="Content.")],
                      speaker_notes=None)],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["notes_quality_valid"] is False


def test_quality_score_perfect(tmp_path):
    good_notes = ("Discuss the key points of this slide in detail. "
                  "Highlight the main finding and explain its implications. "
                  "Transition to the next topic smoothly.")
    ct = SlidesContentTree(
        deck_title="Perfect",
        slides=[Slide(id="s1", slide_type="content", title="Slide",
                      elements=[SlideElement(id="e1", type="body", content="Good content here.")],
                      speaker_notes=good_notes)],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, expected_slide_count=1, content_tree=ct)
    assert result["quality_score"] >= 90


def test_quality_score_degraded(tmp_path):
    ct = SlidesContentTree(
        deck_title="Bad",
        slides=[Slide(id="s1", slide_type="content", title="X",
                      elements=[SlideElement(id="e1", type="body", content="A" * 1000)],
                      speaker_notes=None)],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert result["quality_score"] < 90


def test_small_font_warning(tmp_path):
    """Text with very small font triggers advisory warning."""
    from pptx.util import Pt
    ct = _good_ct()
    path = _export(ct, tmp_path)
    # Manually add a small-font text box
    prs = Presentation(str(path))
    slide = prs.slides[0]
    txBox = slide.shapes.add_textbox(0, 0, 100000, 100000)
    p = txBox.text_frame.paragraphs[0]
    p.text = "tiny text"
    p.font.size = Pt(8)
    prs.save(str(path))
    result = validate_pptx(path)
    assert any("10pt" in w for w in result["layout_warnings"])


def test_result_has_all_keys(tmp_path):
    path = _export(_good_ct(), tmp_path)
    result = validate_pptx(path)
    expected_keys = {
        "valid", "slide_count", "has_notes", "errors",
        "layout_valid", "layout_warnings", "layout_errors",
        "notes_quality_valid", "chart_quality_valid", "quality_score",
    }
    assert expected_keys == set(result.keys())


def test_advisory_long_title_warning(tmp_path):
    ct = SlidesContentTree(
        deck_title="Test",
        slides=[Slide(id="s1", slide_type="content",
                      title="A" * 80,
                      elements=[SlideElement(id="e1", type="body", content="Content here.")],
                      speaker_notes="Good notes for this slide. Discuss the key points and implications.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert any("60 chars" in w for w in result["layout_warnings"])


def test_advisory_excessive_bullets_warning(tmp_path):
    ct = SlidesContentTree(
        deck_title="Test",
        slides=[Slide(id="s1", slide_type="content", title="Bullets",
                      elements=[SlideElement(id="e1", type="bullet_list",
                                             content=["Item"] * 9)],
                      speaker_notes="Walk through each bullet point carefully. Highlight the key takeaways.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert any("9 items" in w for w in result["layout_warnings"])


def test_advisory_sparse_content_warning(tmp_path):
    ct = SlidesContentTree(
        deck_title="Test",
        slides=[Slide(id="s1", slide_type="content", title="Sparse",
                      elements=[SlideElement(id="e1", type="body", content="Hi")],
                      speaker_notes="Discuss this brief slide and explain why minimal content was chosen.")],
    )
    path = _export(ct, tmp_path)
    result = validate_pptx(path, content_tree=ct)
    assert any("sparse" in w.lower() for w in result["layout_warnings"])
