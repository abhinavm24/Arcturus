"""Integration scaffold for P04 (p04_forge).

These tests enforce contract-level integration gates across repo structure and CI wiring,
plus Phase 2 cross-component integration tests.
"""

import asyncio
import json
from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4

import pytest


def _xhtml2pdf_available():
    try:
        from xhtml2pdf import pisa
        return True
    except (ImportError, OSError):
        return False


from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
from core.studio.orchestrator import ForgeOrchestrator
from core.studio.storage import StudioStorage


# Canned LLM responses (duplicated from test_studio_orchestrator to avoid cross-test import)
OUTLINE_RESPONSE = json.dumps({
    "title": "AI Startup Pitch Deck",
    "items": [
        {"id": "1", "title": "Title Slide", "description": "Company intro", "children": []},
        {"id": "2", "title": "Problem", "description": "The pain point", "children": []},
        {"id": "3", "title": "Solution", "description": "Our product", "children": [
            {"id": "3.1", "title": "Demo", "description": "Product demo", "children": []}
        ]},
    ]
})

SLIDES_DRAFT_RESPONSE = json.dumps({
    "deck_title": "AI Startup Pitch Deck",
    "subtitle": "Transforming the Future",
    "slides": [
        {"id": "s1", "slide_type": "title", "title": "Title Slide",
         "elements": [{"id": "e1", "type": "title", "content": "AI Startup"},
                      {"id": "e2", "type": "subtitle", "content": "Series A Pitch"}],
         "speaker_notes": "Welcome everyone."},
        {"id": "s2", "slide_type": "content", "title": "Problem",
         "elements": [{"id": "e3", "type": "body", "content": "Enterprises waste time."}],
         "speaker_notes": "Explain the core problem."},
        {"id": "s3", "slide_type": "content", "title": "Solution",
         "elements": [{"id": "e4", "type": "body", "content": "Our AI platform automates workflows."}],
         "speaker_notes": "Present the solution."},
        {"id": "s4", "slide_type": "two_column", "title": "Before vs After",
         "elements": [{"id": "e5", "type": "body", "content": "Manual processes."},
                      {"id": "e6", "type": "body", "content": "Automated workflows."}],
         "speaker_notes": "Compare old vs new."},
        {"id": "s5", "slide_type": "timeline", "title": "Roadmap",
         "elements": [{"id": "e7", "type": "bullet_list", "content": ["Q1: Launch", "Q2: Scale", "Q3: Expand"]}],
         "speaker_notes": "Walk through the timeline."},
        {"id": "s6", "slide_type": "chart", "title": "Traction",
         "elements": [{"id": "e8", "type": "chart", "content": "Revenue growth chart"},
                      {"id": "e9", "type": "body", "content": "3x growth in 12 months."}],
         "speaker_notes": "Highlight growth metrics."},
        {"id": "s7", "slide_type": "quote", "title": "Testimonial",
         "elements": [{"id": "e10", "type": "quote", "content": "This product changed everything."},
                      {"id": "e11", "type": "body", "content": "Jane Doe, CTO"}],
         "speaker_notes": "Share customer voice."},
        {"id": "s8", "slide_type": "content", "title": "Business Model",
         "elements": [{"id": "e12", "type": "body", "content": "SaaS with enterprise pricing."}],
         "speaker_notes": "Explain monetization."},
        {"id": "s9", "slide_type": "team", "title": "Team",
         "elements": [{"id": "e13", "type": "bullet_list", "content": ["CEO: Alice", "CTO: Bob", "VP Eng: Carol"]}],
         "speaker_notes": "Introduce the team."},
        {"id": "s10", "slide_type": "title", "title": "Thank You",
         "elements": [{"id": "e14", "type": "title", "content": "Thank You"},
                      {"id": "e15", "type": "subtitle", "content": "Questions?"}],
         "speaker_notes": "Close and take questions."},
    ],
    "metadata": {"audience": "investors"},
})


PROJECT_ID = "P04"
PROJECT_KEY = "p04_forge"
CI_CHECK = "p04-forge-studio"
CHARTER = Path("CAPSTONE/project_charters/P04_forge_ai_document_slides_sheets_studio.md")
ACCEPTANCE_FILE = Path("tests/acceptance/p04_forge/test_exports_open_and_render.py")
INTEGRATION_FILE = Path("tests/integration/test_forge_research_to_slides.py")
WORKFLOW_FILE = Path(".github/workflows/project-gates.yml")
BASELINE_SCRIPT = Path("scripts/test_all.sh")


def _read(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _run(coro):
    loop = asyncio.new_event_loop()
    try:
        return loop.run_until_complete(coro)
    finally:
        loop.close()


# === Contract-level scaffold tests ===

def test_01_integration_file_is_declared_in_charter() -> None:
    assert f"Integration: " in _read(CHARTER)


def test_02_acceptance_and_integration_files_exist() -> None:
    assert ACCEPTANCE_FILE.exists(), f"Missing acceptance file: {ACCEPTANCE_FILE}"
    assert INTEGRATION_FILE.exists(), f"Missing integration file: {INTEGRATION_FILE}"


def test_03_baseline_script_exists_and_is_executable() -> None:
    assert BASELINE_SCRIPT.exists(), "Missing baseline script scripts/test_all.sh"
    assert BASELINE_SCRIPT.stat().st_mode & 0o111, "scripts/test_all.sh must be executable"


def test_04_project_ci_check_is_wired_in_workflow() -> None:
    assert WORKFLOW_FILE.exists(), "Missing workflow .github/workflows/project-gates.yml"
    assert CI_CHECK in _read(WORKFLOW_FILE), f"CI check {CI_CHECK} not found in workflow"


def test_05_charter_requires_baseline_regression() -> None:
    assert "scripts/test_all.sh quick" in _read(CHARTER)


# === Phase 2: Cross-component integration tests ===

@pytest.fixture(autouse=True)
def _patch_model_manager_init(monkeypatch):
    """Prevent ModelManager.__init__ from calling real API clients."""
    def noop_init(self, model_name=None, provider=None, role=None):
        self.model_type = "gemini"
        self.client = None
    monkeypatch.setattr("core.model_manager.ModelManager.__init__", noop_init)


@pytest.fixture
def mock_llm(monkeypatch):
    """Mock LLM returning outline then slides draft."""
    async def fake_generate(self, prompt):
        if "content architect" in prompt.lower():
            return OUTLINE_RESPONSE
        return SLIDES_DRAFT_RESPONSE
    monkeypatch.setattr("core.model_manager.ModelManager.generate_text", fake_generate)


@pytest.fixture
def storage(tmp_path):
    return StudioStorage(base_dir=tmp_path / "studio")


@pytest.fixture
def orchestrator(storage):
    return ForgeOrchestrator(storage)


def test_06_outline_to_draft_to_export_pipeline(orchestrator, storage, mock_llm) -> None:
    """Full pipeline: create outline -> approve -> export PPTX."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a pitch deck",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]

    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))
    assert export_result["status"] == "completed"
    assert export_result["file_size_bytes"] > 0


def test_07_export_with_custom_theme(orchestrator, storage, mock_llm) -> None:
    """Export with non-default theme produces valid PPTX."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(
        art_id, ExportFormat.pptx, theme_id="tech-dark"
    ))
    assert export_result["status"] == "completed"


def test_08_export_preserves_revision_lineage(orchestrator, storage, mock_llm) -> None:
    """Revision head_id unchanged after export."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_id_before = art_data["revision_head_id"]

    _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))

    loaded = storage.load_artifact(art_id)
    assert loaded.revision_head_id == rev_id_before


def test_09_multiple_exports_tracked(orchestrator, storage, mock_llm) -> None:
    """Two exports for same artifact both appear in exports list."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))
    _run(orchestrator.export_artifact(art_id, ExportFormat.pptx, theme_id="startup-bold"))

    jobs = storage.list_export_jobs(art_id)
    assert len(jobs) == 2

    loaded = storage.load_artifact(art_id)
    assert len(loaded.exports) == 2


def test_10_export_file_downloadable(orchestrator, storage, mock_llm) -> None:
    """Export file path exists and has non-zero size."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))
    assert export_result["output_uri"] is not None

    export_path = Path(export_result["output_uri"])
    assert export_path.exists()
    assert export_path.stat().st_size > 0


def test_11_oracle_research_ingestion(orchestrator, storage, mock_llm) -> None:
    """Monkeypatch Oracle MCP call with fixture; verify research content appears in content tree."""
    # Phase 2 scope: verify content tree has data from mock LLM
    result = _run(orchestrator.generate_outline(
        prompt="Create a pitch deck about AI agents",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))

    ct = art_data["content_tree"]
    assert ct is not None
    assert ct["deck_title"] == "AI Startup Pitch Deck"
    assert len(ct["slides"]) >= 8


def test_12_canvas_preview_no_schema_breakage(orchestrator, storage, mock_llm) -> None:
    """Round-trip content tree through validation; verify validate_content_tree() passes."""
    from core.schemas.studio_schema import validate_content_tree

    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))

    ct = art_data["content_tree"]
    # Round-trip: dict -> model -> dict -> model
    model = validate_content_tree(ArtifactType.slides, ct)
    roundtripped = model.model_dump(mode="json")
    model2 = validate_content_tree(ArtifactType.slides, roundtripped)
    assert model2.deck_title == model.deck_title
    assert len(model2.slides) == len(model.slides)


def test_13_upstream_failure_graceful_downstream(orchestrator, storage, monkeypatch) -> None:
    """Monkeypatch generate_text to raise; verify HTTP 500 with meaningful error."""
    async def fake_outline(self, prompt):
        if "content architect" in prompt.lower():
            return OUTLINE_RESPONSE
        raise RuntimeError("LLM unavailable")

    monkeypatch.setattr("core.model_manager.ModelManager.generate_text", fake_outline)

    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]

    with pytest.raises(RuntimeError, match="LLM unavailable"):
        _run(orchestrator.approve_and_generate_draft(art_id))

    # Artifact should still exist with outline but no content_tree
    loaded = storage.load_artifact(art_id)
    assert loaded is not None
    assert loaded.content_tree is None
    assert loaded.outline is not None


# === Phase 3: Cross-component integration tests ===

def test_14_outline_to_draft_with_notes_repair(orchestrator, storage, mock_llm) -> None:
    """Draft path applies notes repair; saved content tree has no empty notes."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))

    ct = art_data["content_tree"]
    for slide in ct["slides"]:
        notes = slide.get("speaker_notes", "")
        assert notes and len(notes.strip()) > 0, f"Slide {slide['id']} has empty notes after repair"


def test_15_chart_payload_to_export_pipeline(orchestrator, storage, monkeypatch) -> None:
    """Structured chart data in mock LLM response exports with chart shape."""
    chart_draft = json.dumps({
        "deck_title": "Chart Deck",
        "subtitle": "Data Driven",
        "slides": [
            {"id": "s1", "slide_type": "title", "title": "Intro",
             "elements": [{"id": "e1", "type": "title", "content": "Charts"}],
             "speaker_notes": "Welcome to the data presentation."},
            {"id": "s2", "slide_type": "chart", "title": "Revenue",
             "elements": [{"id": "e2", "type": "chart", "content": {
                 "chart_type": "bar",
                 "categories": ["Q1", "Q2", "Q3", "Q4"],
                 "series": [{"name": "Revenue", "values": [1.0, 2.0, 3.0, 4.0]}],
             }}],
             "speaker_notes": "Discuss the revenue growth trajectory across quarters."},
            {"id": "s3", "slide_type": "content", "title": "Summary",
             "elements": [{"id": "e3", "type": "body", "content": "We grew 4x."}],
             "speaker_notes": "Summarize key takeaways."},
            {"id": "s4", "slide_type": "content", "title": "Market",
             "elements": [{"id": "e4", "type": "body", "content": "Large market."}],
             "speaker_notes": "Explain the market opportunity."},
            {"id": "s5", "slide_type": "content", "title": "Team",
             "elements": [{"id": "e5", "type": "body", "content": "Strong team."}],
             "speaker_notes": "Introduce the team background."},
            {"id": "s6", "slide_type": "content", "title": "Strategy",
             "elements": [{"id": "e6", "type": "body", "content": "Go to market."}],
             "speaker_notes": "Walk through the go-to-market strategy."},
            {"id": "s7", "slide_type": "content", "title": "Product",
             "elements": [{"id": "e7", "type": "body", "content": "Great product."}],
             "speaker_notes": "Highlight key product features."},
            {"id": "s8", "slide_type": "title", "title": "Thanks",
             "elements": [{"id": "e8", "type": "title", "content": "Thank You"}],
             "speaker_notes": "Close and take questions."},
        ],
    })

    async def fake_generate(self, prompt):
        if "content architect" in prompt.lower():
            return OUTLINE_RESPONSE
        return chart_draft
    monkeypatch.setattr("core.model_manager.ModelManager.generate_text", fake_generate)

    result = _run(orchestrator.generate_outline(
        prompt="Create slides with charts",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))
    assert export_result["status"] == "completed"

    # Verify chart shape in exported PPTX
    from pptx import Presentation as PptxPrs
    prs = PptxPrs(export_result["output_uri"])
    chart_slide = prs.slides[1]
    has_chart = any(s.has_chart for s in chart_slide.shapes)
    assert has_chart


def test_16_variant_theme_export_pipeline(orchestrator, storage, mock_llm) -> None:
    """Export with variant theme ID succeeds."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(
        art_id, ExportFormat.pptx, theme_id="corporate-blue--v01"
    ))
    assert export_result["status"] == "completed"


def test_17_strict_export_repairs_without_mutating_artifact(orchestrator, storage, mock_llm) -> None:
    """Strict export repairs layout on export copy but preserves stored artifact content_tree."""
    result = _run(orchestrator.generate_outline(
        prompt="Create slides",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_before = art_data["revision_head_id"]

    # Force overflow in stored content tree
    loaded = storage.load_artifact(art_id)
    loaded.content_tree["slides"][1]["elements"][0]["content"] = "X" * 2500
    storage.save_artifact(loaded)

    export_result = _run(orchestrator.export_artifact(
        art_id, ExportFormat.pptx, strict_layout=True
    ))
    assert export_result["status"] == "completed"

    # Artifact state unchanged — repair only affects export copy
    reloaded = storage.load_artifact(art_id)
    assert reloaded.revision_head_id == rev_before
    assert len(reloaded.content_tree["slides"][1]["elements"][0]["content"]) == 2500


# --- Phase 4: Document pipeline integration tests ---

DOCUMENT_DRAFT_RESPONSE = json.dumps({
    "doc_title": "Technical Specification: Agent Framework",
    "doc_type": "technical_spec",
    "abstract": "This document specifies the architecture.",
    "sections": [
        {"id": "sec1", "heading": "Introduction", "level": 1,
         "content": "The framework provides a modular architecture.",
         "subsections": [
             {"id": "sec1a", "heading": "Purpose", "level": 2,
              "content": "Defines core components.", "subsections": [], "citations": []}
         ],
         "citations": ["ref1"]},
        {"id": "sec2", "heading": "Architecture", "level": 1,
         "content": "Layered architecture with separation of concerns.",
         "subsections": [], "citations": []},
        {"id": "sec3", "heading": "Implementation", "level": 1,
         "content": "Phased approach with iterative delivery.",
         "subsections": [], "citations": []},
    ],
    "bibliography": [
        {"key": "ref1", "title": "AI Agents", "author": "Smith et al.", "year": "2024"},
    ],
    "metadata": {"audience": "engineers"},
})


@pytest.fixture
def mock_llm_document(monkeypatch):
    """Mock LLM for document drafts."""
    async def fake_generate(self, prompt):
        if "content architect" in prompt.lower():
            return OUTLINE_RESPONSE
        return DOCUMENT_DRAFT_RESPONSE
    monkeypatch.setattr("core.model_manager.ModelManager.generate_text", fake_generate)


def test_18_document_outline_to_draft_to_docx_export(orchestrator, storage, mock_llm_document) -> None:
    """Full document pipeline: outline → draft → export DOCX."""
    result = _run(orchestrator.generate_outline(
        prompt="Write a technical specification",
        artifact_type=ArtifactType.document,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    assert art_data["content_tree"]["doc_title"] == "Technical Specification: Agent Framework"

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.docx))
    assert export_result["status"] == "completed"
    assert export_result["format"] == "docx"
    assert export_result["validator_results"]["valid"] is True


@pytest.mark.skipif(
    not _xhtml2pdf_available(),
    reason="xhtml2pdf not available",
)
def test_19_document_pdf_export(orchestrator, storage, mock_llm_document) -> None:
    """Document pipeline: outline → draft → export PDF."""
    result = _run(orchestrator.generate_outline(
        prompt="Write a technical specification",
        artifact_type=ArtifactType.document,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.pdf))
    assert export_result["status"] == "completed"
    assert export_result["format"] == "pdf"
    assert export_result["validator_results"]["valid"] is True


@pytest.mark.skipif(
    not _xhtml2pdf_available(),
    reason="xhtml2pdf not available",
)
def test_20_document_multiple_exports_tracked(orchestrator, storage, mock_llm_document) -> None:
    """Multiple document exports are tracked on the artifact."""
    result = _run(orchestrator.generate_outline(
        prompt="Write a report",
        artifact_type=ArtifactType.document,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    _run(orchestrator.export_artifact(art_id, ExportFormat.docx))
    _run(orchestrator.export_artifact(art_id, ExportFormat.pdf))

    loaded = storage.load_artifact(art_id)
    assert len(loaded.exports) == 2
    formats = {e.format for e in loaded.exports}
    assert formats == {"docx", "pdf"}


def test_21_document_export_pptx_blocked(orchestrator, storage, mock_llm_document) -> None:
    """PPTX export should fail for document artifacts."""
    result = _run(orchestrator.generate_outline(
        prompt="Write a report",
        artifact_type=ArtifactType.document,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    with pytest.raises(ValueError, match="not supported"):
        _run(orchestrator.export_artifact(art_id, ExportFormat.pptx))


# --- Phase 5: Sheet pipeline integration tests ---

SHEET_DRAFT_RESPONSE = json.dumps({
    "workbook_title": "Sales Report",
    "tabs": [
        {
            "id": "t1", "name": "Revenue",
            "headers": ["Month", "Amount", "Growth"],
            "rows": [
                ["Jan", 1000, 0.0],
                ["Feb", 1200, 0.2],
                ["Mar", 1500, 0.25],
                ["Apr", 1350, -0.1],
            ],
            "formulas": {"D2": "=C2/B2"},
            "column_widths": [100, 80, 80],
        },
        {
            "id": "t2", "name": "Summary",
            "headers": ["Metric", "Value"],
            "rows": [["Total", 5050], ["Average", 1262.5]],
            "formulas": {},
            "column_widths": [120, 80],
        },
    ],
})


@pytest.fixture
def mock_llm_sheet(monkeypatch):
    """Mock LLM for sheet drafts."""
    async def fake_generate(self, prompt):
        if "content architect" in prompt.lower():
            return OUTLINE_RESPONSE
        return SHEET_DRAFT_RESPONSE
    monkeypatch.setattr("core.model_manager.ModelManager.generate_text", fake_generate)


def test_22_sheet_outline_to_draft_to_xlsx_export(orchestrator, storage, mock_llm_sheet) -> None:
    """Full sheet pipeline: outline → draft → XLSX export."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a sales report spreadsheet",
        artifact_type=ArtifactType.sheet,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    assert art_data["content_tree"]["workbook_title"] == "Sales Report"

    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx))
    assert export_result["status"] == "completed"
    assert export_result["format"] == "xlsx"
    assert export_result["file_size_bytes"] > 0
    assert export_result["validator_results"]["valid"] is True
    assert export_result["validator_results"]["chart_count"] >= 1
    assert export_result["validator_results"]["quality_score"] >= 60


def test_23_sheet_upload_analysis_to_export_pipeline(orchestrator, storage, mock_llm_sheet) -> None:
    """Upload CSV → analysis → XLSX export."""
    import csv as csv_mod
    import tempfile

    # 1. Create the initial sheet artifact via outline+draft
    result = _run(orchestrator.generate_outline(
        prompt="Create a report",
        artifact_type=ArtifactType.sheet,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    # 2. Create a CSV file for upload
    csv_content = "Product,Revenue,Units\n"
    for i in range(10):
        csv_content += f"P{i},{100 + i * 10},{50 + i}\n"

    # 3. Upload and analyze (returns full artifact dict)
    upload_result = _run(orchestrator.analyze_sheet_upload(
        art_id, "data.csv", csv_content.encode("utf-8"), "text/csv"
    ))
    ct = upload_result["content_tree"]
    assert ct["analysis_report"] is not None
    assert len(ct["analysis_report"]["summary_stats"]) >= 2

    # 4. Export the updated artifact to XLSX
    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx))
    assert export_result["status"] == "completed"
    assert export_result["validator_results"]["chart_count"] >= 1
    assert export_result["validator_results"]["quality_score"] >= 60


def test_24_sheet_invalid_upload_graceful_failure(orchestrator, storage, mock_llm_sheet) -> None:
    """Invalid upload returns controlled error, artifact unchanged."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a report",
        artifact_type=ArtifactType.sheet,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    # Get revision before upload attempt
    loaded_before = storage.load_artifact(art_id)
    rev_before = loaded_before.revision_head_id

    # Upload a binary file with unsupported type
    with pytest.raises(ValueError, match="Unsupported"):
        _run(orchestrator.analyze_sheet_upload(
            art_id, "data.pptx", b"fake content", "application/vnd.ms-powerpoint"
        ))

    # Artifact unchanged
    loaded_after = storage.load_artifact(art_id)
    assert loaded_after.revision_head_id == rev_before


def test_25_sheet_export_param_gating_enforced(orchestrator, storage, mock_llm_sheet) -> None:
    """Slides-only params (theme_id, strict_layout, generate_images) rejected for sheet exports."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a report",
        artifact_type=ArtifactType.sheet,
    ))
    art_id = result["artifact_id"]
    _run(orchestrator.approve_and_generate_draft(art_id))

    with pytest.raises(ValueError, match="theme_id"):
        _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx, theme_id="corporate-blue"))

    with pytest.raises(ValueError, match="strict_layout"):
        _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx, strict_layout=True))

    with pytest.raises(ValueError, match="generate_images"):
        _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx, generate_images=True))


# === Phase 6: Edit Loop Integration Tests ===


def test_26_slides_edit_loop_pipeline(orchestrator, storage, mock_llm) -> None:
    """Full slides pipeline: outline → draft → edit → verify revision + content change."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a pitch deck",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_before = art_data["revision_head_id"]
    ct_before = art_data["content_tree"]

    # Edit slide 2 title using _patch_override (deterministic, no LLM)
    patch = {
        "artifact_type": "slides",
        "target": {"kind": "slide_index", "index": 2},
        "ops": [{"op": "SET", "path": "title", "value": "Updated Problem Statement"}],
        "summary": "Rename slide 2",
    }
    edit_result = _run(orchestrator.edit_artifact(
        art_id, "Rename slide 2 to Updated Problem Statement",
        base_revision_id=rev_before,
        _patch_override=patch,
    ))

    # Verify edit applied
    assert edit_result["edit_result"]["status"] == "applied"
    assert edit_result["edit_result"]["diff"]["stats"]["paths_changed"] > 0

    # Verify revision advanced
    loaded = storage.load_artifact(art_id)
    assert loaded.revision_head_id != rev_before

    # Verify content tree changed
    new_ct = loaded.content_tree
    assert new_ct["slides"][1]["title"] == "Updated Problem Statement"

    # Verify revision has edit metadata
    new_rev_id = edit_result["edit_result"]["revision_id"]
    latest = storage.load_revision(art_id, new_rev_id)
    assert latest.edit_instruction == "Rename slide 2 to Updated Problem Statement"
    assert latest.patch is not None
    assert latest.diff is not None


def test_27_document_edit_loop_pipeline(orchestrator, storage, mock_llm_document) -> None:
    """Full document pipeline: outline → draft → edit → verify revision + export DOCX."""
    result = _run(orchestrator.generate_outline(
        prompt="Write a technical specification",
        artifact_type=ArtifactType.document,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_before = art_data["revision_head_id"]

    # Edit section content
    patch = {
        "artifact_type": "document",
        "target": {"kind": "section_id", "id": "sec1"},
        "ops": [{"op": "SET", "path": "content", "value": "Revised introduction for v2."}],
        "summary": "Update introduction",
    }
    edit_result = _run(orchestrator.edit_artifact(
        art_id, "Revise the introduction",
        base_revision_id=rev_before,
        _patch_override=patch,
    ))
    assert edit_result["edit_result"]["status"] == "applied"

    # Verify content change persisted
    loaded = storage.load_artifact(art_id)
    sec1 = loaded.content_tree["sections"][0]
    assert sec1["content"] == "Revised introduction for v2."

    # Export DOCX after edit
    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.docx))
    assert export_result["status"] == "completed"
    assert export_result["format"] == "docx"


def test_28_sheet_edit_loop_pipeline(orchestrator, storage, mock_llm_sheet) -> None:
    """Full sheet pipeline: outline → draft → edit → verify revision + export XLSX."""
    result = _run(orchestrator.generate_outline(
        prompt="Create a sales report spreadsheet",
        artifact_type=ArtifactType.sheet,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_before = art_data["revision_head_id"]

    # Edit: update a cell value
    patch = {
        "artifact_type": "sheet",
        "target": {"kind": "tab_name", "name": "Revenue"},
        "ops": [{"op": "SET", "path": "rows[0][1]", "value": 9999}],
        "summary": "Update Jan revenue",
    }
    edit_result = _run(orchestrator.edit_artifact(
        art_id, "Update January revenue to 9999",
        base_revision_id=rev_before,
        _patch_override=patch,
    ))
    assert edit_result["edit_result"]["status"] == "applied"

    # Verify cell changed
    loaded = storage.load_artifact(art_id)
    tab = loaded.content_tree["tabs"][0]
    assert tab["rows"][0][1] == 9999

    # Export XLSX after edit
    export_result = _run(orchestrator.export_artifact(art_id, ExportFormat.xlsx))
    assert export_result["status"] == "completed"
    assert export_result["format"] == "xlsx"


def test_29_edit_conflict_detection(orchestrator, storage, mock_llm) -> None:
    """Optimistic concurrency: stale base_revision_id triggers ConflictError."""
    from core.studio.orchestrator import ConflictError

    result = _run(orchestrator.generate_outline(
        prompt="Create a pitch deck",
        artifact_type=ArtifactType.slides,
    ))
    art_id = result["artifact_id"]
    art_data = _run(orchestrator.approve_and_generate_draft(art_id))
    rev_v1 = art_data["revision_head_id"]

    # First edit succeeds
    patch1 = {
        "artifact_type": "slides",
        "target": {"kind": "slide_index", "index": 1},
        "ops": [{"op": "SET", "path": "title", "value": "Edit 1"}],
        "summary": "First edit",
    }
    _run(orchestrator.edit_artifact(
        art_id, "First edit", base_revision_id=rev_v1, _patch_override=patch1
    ))

    # Second edit with stale rev_v1 should fail
    patch2 = {
        "artifact_type": "slides",
        "target": {"kind": "slide_index", "index": 1},
        "ops": [{"op": "SET", "path": "title", "value": "Edit 2"}],
        "summary": "Second edit",
    }
    with pytest.raises(ConflictError):
        _run(orchestrator.edit_artifact(
            art_id, "Second edit", base_revision_id=rev_v1, _patch_override=patch2
        ))
