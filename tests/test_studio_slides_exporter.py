"""Tests for core/studio/slides/exporter.py — PPTX rendering."""

import pytest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
from core.studio.slides.exporter import (
    SLIDE_HEIGHT, SLIDE_WIDTH, _DESIGN_TOKENS, _blend_color, _build_chart_palette,
    export_to_pptx,
)
from core.studio.slides.themes import get_theme, generate_theme_variant


@pytest.fixture
def sample_content_tree():
    return SlidesContentTree(
        deck_title="Test Deck",
        subtitle="Test Subtitle",
        slides=[
            Slide(
                id="s1",
                slide_type="title",
                title="Welcome",
                elements=[
                    SlideElement(id="e1", type="title", content="Welcome"),
                    SlideElement(id="e2", type="subtitle", content="A test deck"),
                ],
                speaker_notes="Open with a greeting.",
            ),
            Slide(
                id="s2",
                slide_type="content",
                title="Agenda",
                elements=[
                    SlideElement(id="e3", type="title", content="Agenda"),
                    SlideElement(id="e4", type="bullet_list", content=["Item 1", "Item 2", "Item 3"]),
                ],
                speaker_notes="Walk through the agenda items.",
            ),
            Slide(
                id="s3",
                slide_type="two_column",
                title="Comparison",
                elements=[
                    SlideElement(id="e5", type="body", content="Left column content."),
                    SlideElement(id="e6", type="body", content="Right column content."),
                ],
                speaker_notes="Compare both sides.",
            ),
            Slide(
                id="s4",
                slide_type="quote",
                title="Inspiration",
                elements=[
                    SlideElement(id="e7", type="quote", content="The best way to predict the future is to invent it."),
                    SlideElement(id="e8", type="body", content="Alan Kay"),
                ],
                speaker_notes="Pause for effect.",
            ),
            Slide(
                id="s5",
                slide_type="code",
                title="Code Example",
                elements=[
                    SlideElement(id="e9", type="code", content="def hello():\n    print('Hello world')"),
                ],
                speaker_notes="Walk through the code.",
            ),
        ],
    )


@pytest.fixture
def default_theme():
    return get_theme()


def test_export_creates_file(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    result = export_to_pptx(sample_content_tree, default_theme, output)
    assert result == output
    assert output.exists()
    assert output.stat().st_size > 0


def test_export_slide_count(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 5


def test_export_speaker_notes_present(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    has_notes = False
    for slide in prs.slides:
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                has_notes = True
                break
        except Exception:
            continue
    assert has_notes


def test_export_all_slides_have_notes(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    for i, slide in enumerate(prs.slides):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        assert notes, f"Slide {i+1} missing speaker notes"


def test_export_slide_dimensions(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    assert prs.slide_width == SLIDE_WIDTH
    assert prs.slide_height == SLIDE_HEIGHT


def test_export_title_slide_rendering(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    title_slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in title_slide.shapes if shape.has_text_frame]
    assert any("Welcome" in t for t in texts)


def test_export_content_slide_rendering(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    content_slide = prs.slides[1]
    texts = [shape.text_frame.text for shape in content_slide.shapes if shape.has_text_frame]
    all_text = " ".join(texts)
    assert "Agenda" in all_text
    assert "Item 1" in all_text


def test_export_two_column_rendering(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    two_col_slide = prs.slides[2]
    # Two-column should have at least 3 shapes (title + 2 columns)
    assert len(two_col_slide.shapes) >= 3


def test_export_quote_slide_rendering(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    quote_slide = prs.slides[3]
    texts = [shape.text_frame.text for shape in quote_slide.shapes if shape.has_text_frame]
    all_text = " ".join(texts)
    assert "predict the future" in all_text


def test_export_code_slide_rendering(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    code_slide = prs.slides[4]
    texts = [shape.text_frame.text for shape in code_slide.shapes if shape.has_text_frame]
    all_text = " ".join(texts)
    assert "hello" in all_text


def test_export_with_different_themes(sample_content_tree, tmp_path):
    for theme_id in ["corporate-blue", "tech-dark", "startup-bold"]:
        theme = get_theme(theme_id)
        output = tmp_path / f"{theme_id}.pptx"
        export_to_pptx(sample_content_tree, theme, output)
        prs = Presentation(str(output))
        assert len(prs.slides) == 5


def test_export_rejects_empty_slides_list(default_theme, tmp_path):
    ct = SlidesContentTree(deck_title="Empty", slides=[])
    output = tmp_path / "empty.pptx"
    # Empty slides produces a valid but empty PPTX (no crash)
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 0


def test_export_output_directory_created(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "nested" / "deep" / "test.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    assert output.exists()


def test_export_unknown_slide_type_fallback(default_theme, tmp_path):
    ct = SlidesContentTree(
        deck_title="Unknown Types",
        slides=[
            Slide(
                id="s1",
                slide_type="mystery_type",
                title="Unknown",
                elements=[
                    SlideElement(id="e1", type="body", content="Falls back to content renderer."),
                ],
                speaker_notes="Should not crash.",
            ),
        ],
    )
    output = tmp_path / "unknown.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 1


# === Phase 3: Chart, Chrome, Card, and Visual Tests ===

def _bar_chart_spec():
    return {
        "chart_type": "bar",
        "title": "Revenue",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{"name": "Revenue", "values": [1.2, 1.8, 2.6, 3.1]}],
    }


def _chart_ct(chart_content):
    return SlidesContentTree(
        deck_title="Chart Test",
        slides=[
            Slide(id="s1", slide_type="title", title="Title",
                  elements=[SlideElement(id="e0", type="title", content="Title")],
                  speaker_notes="Opening."),
            Slide(id="s2", slide_type="chart", title="Chart Slide",
                  elements=[SlideElement(id="e1", type="chart", content=chart_content)],
                  speaker_notes="Chart notes."),
            Slide(id="s3", slide_type="content", title="End",
                  elements=[SlideElement(id="e2", type="body", content="Closing.")],
                  speaker_notes="Closing notes."),
        ],
    )


def test_export_chart_bar_native(default_theme, tmp_path):
    ct = _chart_ct(_bar_chart_spec())
    output = tmp_path / "bar.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    chart_slide = prs.slides[1]
    has_chart = any(shape.has_chart for shape in chart_slide.shapes)
    assert has_chart, "Bar chart slide should contain a chart shape"


def test_export_chart_line_native(default_theme, tmp_path):
    spec = {"chart_type": "line", "categories": ["Jan", "Feb", "Mar"],
            "series": [{"name": "Growth", "values": [10, 20, 30]}]}
    ct = _chart_ct(spec)
    output = tmp_path / "line.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert any(s.has_chart for s in prs.slides[1].shapes)


def test_export_chart_pie_native(default_theme, tmp_path):
    spec = {"chart_type": "pie", "categories": ["A", "B", "C"],
            "series": [{"name": "Share", "values": [40, 35, 25]}]}
    ct = _chart_ct(spec)
    output = tmp_path / "pie.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert any(s.has_chart for s in prs.slides[1].shapes)


def test_export_chart_scatter_native(default_theme, tmp_path):
    spec = {"chart_type": "scatter",
            "points": [{"x": 1.0, "y": 2.0}, {"x": 3.0, "y": 4.0}]}
    ct = _chart_ct(spec)
    output = tmp_path / "scatter.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert any(s.has_chart for s in prs.slides[1].shapes)


def test_export_chart_funnel_fallback(default_theme, tmp_path):
    spec = {"chart_type": "funnel", "categories": ["Top", "Mid", "Bot"],
            "series": [{"name": "F", "values": [100, 60, 20]}]}
    ct = _chart_ct(spec)
    output = tmp_path / "funnel.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert any(s.has_chart for s in prs.slides[1].shapes)


def test_export_chart_invalid_fallback(default_theme, tmp_path):
    ct = _chart_ct({"invalid": True})
    output = tmp_path / "invalid.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 3  # No crash


def test_export_chart_string_fallback(default_theme, tmp_path):
    ct = _chart_ct("Revenue growth chart")
    output = tmp_path / "string.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    chart_slide = prs.slides[1]
    texts = [s.text_frame.text for s in chart_slide.shapes if s.has_text_frame]
    assert any("[Chart:" in t for t in texts)


def test_export_chart_pie_multi_series(default_theme, tmp_path):
    spec = {"chart_type": "pie", "categories": ["A", "B"],
            "series": [
                {"name": "S1", "values": [60, 40]},
                {"name": "S2", "values": [50, 50]},
                {"name": "S3", "values": [30, 70]},
            ]}
    ct = _chart_ct(spec)
    output = tmp_path / "pie_multi.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    assert any(s.has_chart for s in prs.slides[1].shapes)


def test_export_gradient_background(tmp_path):
    theme = generate_theme_variant("corporate-blue", 2)  # seed 2 gets gradient
    ct = SlidesContentTree(
        deck_title="Gradient",
        slides=[Slide(id="s1", slide_type="content", title="T",
                      elements=[SlideElement(id="e1", type="body", content="Content.")],
                      speaker_notes="Notes.")],
    )
    output = tmp_path / "gradient.pptx"
    export_to_pptx(ct, theme, output)
    assert output.exists()
    assert output.stat().st_size > 0


def test_export_with_variant_theme(tmp_path):
    theme = get_theme("corporate-blue--v01")
    ct = SlidesContentTree(
        deck_title="Variant",
        slides=[Slide(id="s1", slide_type="content", title="T",
                      elements=[SlideElement(id="e1", type="body", content="Text.")],
                      speaker_notes="Notes.")],
    )
    output = tmp_path / "variant.pptx"
    export_to_pptx(ct, theme, output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 1


def test_export_text_frame_has_margins(default_theme, tmp_path):
    ct = SlidesContentTree(
        deck_title="Margins",
        slides=[Slide(id="s1", slide_type="content", title="T",
                      elements=[SlideElement(id="e1", type="body", content="Text content.")],
                      speaker_notes="Notes.")],
    )
    output = tmp_path / "margins.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            tf = shape.text_frame
            assert tf.margin_left is not None
            assert tf.margin_left > 0
            break


def test_export_content_slide_has_card_shape(default_theme, tmp_path):
    ct = SlidesContentTree(
        deck_title="Card",
        slides=[Slide(id="s1", slide_type="content", title="T",
                      elements=[SlideElement(id="e1", type="body", content="Body text.")],
                      speaker_notes="Notes.")],
    )
    output = tmp_path / "card.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    slide = prs.slides[0]
    # Content slide should have rectangle shapes (card + accent strip)
    rect_count = sum(1 for s in slide.shapes if hasattr(s, 'shape_type') and s.shape_type == 1)
    assert rect_count >= 1, "Content slide should have card rectangle shapes"


def test_export_comparison_has_two_cards(default_theme, tmp_path):
    ct = SlidesContentTree(
        deck_title="Compare",
        slides=[Slide(id="s1", slide_type="comparison", title="Compare",
                      elements=[
                          SlideElement(id="e1", type="body", content="Left."),
                          SlideElement(id="e2", type="body", content="Right."),
                      ],
                      speaker_notes="Compare both sides.")],
    )
    output = tmp_path / "compare.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    slide = prs.slides[0]
    rect_count = sum(1 for s in slide.shapes if hasattr(s, 'shape_type') and s.shape_type == 1)
    assert rect_count >= 2, "Comparison slide should have at least 2 card rectangles"


def test_export_body_slide_has_accent_bar(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "chrome.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    # Check middle slide (index 1) for accent bar
    slide = prs.slides[1]
    has_full_width_rect = False
    for shape in slide.shapes:
        if hasattr(shape, 'shape_type') and shape.shape_type == 1:
            if shape.width >= SLIDE_WIDTH - Inches(0.1):
                has_full_width_rect = True
                break
    assert has_full_width_rect, "Content slide should have full-width accent bar"


def test_export_body_slide_has_slide_number(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "numbers.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    slide = prs.slides[1]
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    has_number = any("/" in t for t in texts)
    assert has_number, "Content slide should have slide number"


def test_export_title_slide_no_chrome(sample_content_tree, default_theme, tmp_path):
    output = tmp_path / "no_chrome.pptx"
    export_to_pptx(sample_content_tree, default_theme, output)
    prs = Presentation(str(output))
    title_slide = prs.slides[0]
    texts = [s.text_frame.text for s in title_slide.shapes if s.has_text_frame]
    assert not any("/" in t and t.strip().replace(" ", "").replace("/", "").isdigit()
                    for t in texts), "Title slide should not have slide number"


def test_export_chart_has_styled_series(default_theme, tmp_path):
    ct = _chart_ct(_bar_chart_spec())
    output = tmp_path / "styled.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    chart_slide = prs.slides[1]
    for shape in chart_slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            assert len(chart.series) >= 1
            return
    pytest.fail("No chart shape found")


def test_export_chart_has_gridlines(default_theme, tmp_path):
    ct = _chart_ct(_bar_chart_spec())
    output = tmp_path / "gridlines.pptx"
    export_to_pptx(ct, default_theme, output)
    prs = Presentation(str(output))
    chart_slide = prs.slides[1]
    for shape in chart_slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            assert chart.value_axis.has_major_gridlines
            return
    pytest.fail("No chart shape found")
