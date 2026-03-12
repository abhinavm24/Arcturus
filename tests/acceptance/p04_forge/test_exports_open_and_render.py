"""Acceptance scaffold for P04 (p04_forge).

Replace these contract tests with feature-level assertions as implementation matures.
"""

from pathlib import Path

import pytest

def _xhtml2pdf_available():
    try:
        from xhtml2pdf import pisa
        return True
    except (ImportError, OSError):
        return False


PROJECT_ID = "P04"
PROJECT_KEY = "p04_forge"
CI_CHECK = "p04-forge-studio"
CHARTER = Path("CAPSTONE/project_charters/P04_forge_ai_document_slides_sheets_studio.md")
DELIVERY_README = Path("CAPSTONE/project_charters/P04_DELIVERY_README.md")
DEMO_SCRIPT = Path("scripts/demos/p04_forge.sh")
THIS_FILE = Path("tests/acceptance/p04_forge/test_exports_open_and_render.py")


def _sample_content_tree_dict() -> dict:
    """Return a 10-slide content tree dict for testing."""
    return {
        "deck_title": "AI Startup Pitch Deck",
        "subtitle": "Transforming the Future",
        "slides": [
            {"id": "s1", "slide_type": "title", "title": "Title Slide",
             "elements": [{"id": "e1", "type": "title", "content": "AI Startup"},
                          {"id": "e2", "type": "subtitle", "content": "Series A Pitch"}],
             "speaker_notes": "Welcome everyone."},
            {"id": "s2", "slide_type": "content", "title": "Problem",
             "elements": [{"id": "e3", "type": "body", "content": "Enterprises waste time."}],
             "speaker_notes": "Explain the core problem."},
            {"id": "s3", "slide_type": "content", "title": "Solution",
             "elements": [{"id": "e4", "type": "body", "content": "AI platform."}],
             "speaker_notes": "Present the solution."},
            {"id": "s4", "slide_type": "two_column", "title": "Before vs After",
             "elements": [{"id": "e5", "type": "body", "content": "Manual."},
                          {"id": "e6", "type": "body", "content": "Automated."}],
             "speaker_notes": "Compare."},
            {"id": "s5", "slide_type": "timeline", "title": "Roadmap",
             "elements": [{"id": "e7", "type": "bullet_list", "content": ["Q1", "Q2", "Q3"]}],
             "speaker_notes": "Timeline."},
            {"id": "s6", "slide_type": "chart", "title": "Traction",
             "elements": [{"id": "e8", "type": "chart", "content": "Growth chart"}],
             "speaker_notes": "Metrics."},
            {"id": "s7", "slide_type": "quote", "title": "Testimonial",
             "elements": [{"id": "e10", "type": "quote", "content": "Amazing."}],
             "speaker_notes": "Customer."},
            {"id": "s8", "slide_type": "content", "title": "Business Model",
             "elements": [{"id": "e12", "type": "body", "content": "SaaS."}],
             "speaker_notes": "Monetization."},
            {"id": "s9", "slide_type": "team", "title": "Team",
             "elements": [{"id": "e13", "type": "bullet_list", "content": ["Alice", "Bob"]}],
             "speaker_notes": "Team."},
            {"id": "s10", "slide_type": "title", "title": "Thank You",
             "elements": [{"id": "e14", "type": "title", "content": "Thank You"}],
             "speaker_notes": "Close."},
        ],
        "metadata": {"audience": "investors"},
    }


def _charter_text() -> str:
    return CHARTER.read_text(encoding="utf-8")


def test_01_charter_exists() -> None:
    assert CHARTER.exists(), f"Missing charter: {CHARTER}"


def test_02_expanded_gate_contract_present() -> None:
    assert "Expanded Mandatory Test Gate Contract (10 Hard Conditions)" in _charter_text()


def test_03_acceptance_path_declared_in_charter() -> None:
    assert f"Acceptance: " in _charter_text()


def test_04_demo_script_exists() -> None:
    assert DEMO_SCRIPT.exists(), f"Missing demo script: {DEMO_SCRIPT}"


def test_05_demo_script_is_executable() -> None:
    assert DEMO_SCRIPT.stat().st_mode & 0o111, f"Demo script not executable: {DEMO_SCRIPT}"


def test_06_delivery_readme_exists() -> None:
    assert DELIVERY_README.exists(), f"Missing delivery README: {DELIVERY_README}"


def test_07_delivery_readme_has_required_sections() -> None:
    required = [
        "## 1. Scope Delivered",
        "## 2. Architecture Changes",
        "## 3. API And UI Changes",
        "## 4. Mandatory Test Gate Definition",
        "## 5. Test Evidence",
        "## 8. Known Gaps",
        "## 10. Demo Steps",
    ]
    text = DELIVERY_README.read_text(encoding="utf-8")
    for section in required:
        assert section in text, f"Missing section {section} in {DELIVERY_README}"


def test_08_ci_check_declared_in_charter() -> None:
    assert f"CI required check: " in _charter_text()


# === Phase 2: Export + Render Functional Tests ===

def test_09_slides_content_tree_validates() -> None:
    """A sample SlidesContentTree passes Pydantic validation."""
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree

    ct = SlidesContentTree(
        deck_title="Test Deck",
        slides=[
            Slide(
                id="s1",
                slide_type="title",
                title="Welcome",
                elements=[SlideElement(id="e1", type="title", content="Hello")],
                speaker_notes="Greet the audience.",
            ),
        ],
    )
    assert ct.deck_title == "Test Deck"
    assert len(ct.slides) == 1


def test_10_pptx_export_produces_file(tmp_path) -> None:
    """export_to_pptx() creates a valid PPTX file."""
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.exporter import export_to_pptx
    from core.studio.slides.themes import get_theme

    ct = SlidesContentTree(
        deck_title="Export Test",
        slides=[
            Slide(
                id="s1", slide_type="title", title="Hello",
                elements=[SlideElement(id="e1", type="title", content="Hello")],
                speaker_notes="Open.",
            ),
        ],
    )
    output = tmp_path / "test.pptx"
    export_to_pptx(ct, get_theme(), output)
    assert output.exists()
    assert output.stat().st_size > 0


def test_11_pptx_open_validation_passes(tmp_path) -> None:
    """validate_pptx() confirms the exported file opens cleanly."""
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.exporter import export_to_pptx
    from core.studio.slides.themes import get_theme
    from core.studio.slides.validator import validate_pptx

    ct = SlidesContentTree(
        deck_title="Validation Test",
        slides=[
            Slide(
                id="s1", slide_type="content", title="Slide",
                elements=[SlideElement(id="e1", type="body", content="Content here.")],
                speaker_notes="Talk about this.",
            ),
        ],
    )
    output = tmp_path / "valid.pptx"
    export_to_pptx(ct, get_theme(), output)

    result = validate_pptx(output, expected_slide_count=1)
    assert result["valid"] is True
    assert result["slide_count"] == 1


def test_12_slide_count_in_range() -> None:
    """clamp_slide_count enforces [MIN_SLIDES, MAX_SLIDES] range."""
    from core.studio.slides.generator import clamp_slide_count, MIN_SLIDES, MAX_SLIDES

    assert clamp_slide_count(1) == MIN_SLIDES
    assert clamp_slide_count(50) == MAX_SLIDES
    assert MIN_SLIDES <= clamp_slide_count(10) <= MAX_SLIDES


def test_13_speaker_notes_present_in_export(tmp_path) -> None:
    """At least one slide has speaker notes in the PPTX."""
    from pptx import Presentation as PptxPresentation
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.exporter import export_to_pptx
    from core.studio.slides.themes import get_theme

    ct = SlidesContentTree(
        deck_title="Notes Test",
        slides=[
            Slide(
                id="s1", slide_type="title", title="Intro",
                elements=[SlideElement(id="e1", type="title", content="Intro")],
                speaker_notes="This is a speaker note.",
            ),
        ],
    )
    output = tmp_path / "notes.pptx"
    export_to_pptx(ct, get_theme(), output)

    prs = PptxPresentation(str(output))
    has_notes = any(
        s.notes_slide.notes_text_frame.text.strip()
        for s in prs.slides
    )
    assert has_notes


def test_14_export_job_status_completed(tmp_path) -> None:
    """Export job ends with status=completed for valid input."""
    import asyncio
    from core.schemas.studio_schema import ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage
    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    import json
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType

    art_id = str(uuid4())
    ct = _sample_content_tree_dict()
    artifact = Artifact(
        id=art_id,
        type=ArtifactType.slides,
        title="Test",
        created_at=datetime.now(timezone.utc),
        updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.pptx)
        )
    finally:
        loop.close()

    assert result["status"] == "completed"


def test_15_invalid_format_returns_error() -> None:
    """Requesting unsupported format returns controlled error."""
    from core.schemas.studio_schema import ExportFormat
    with pytest.raises(ValueError):
        ExportFormat("odt")


def test_16_export_job_has_validator_results(tmp_path) -> None:
    """Export job includes validator_results."""
    import asyncio
    import json
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage
    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    art_id = str(uuid4())
    ct = _sample_content_tree_dict()
    artifact = Artifact(
        id=art_id,
        type=ArtifactType.slides,
        title="Test",
        created_at=datetime.now(timezone.utc),
        updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.pptx)
        )
    finally:
        loop.close()

    vr = result["validator_results"]
    assert "valid" in vr
    assert "slide_count" in vr
    assert "has_notes" in vr
    assert "errors" in vr


def test_17_layout_validator_detects_overflow(tmp_path) -> None:
    """Slide with 2000+ chars triggers layout_valid=False."""
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.exporter import export_to_pptx
    from core.studio.slides.themes import get_theme
    from core.studio.slides.validator import validate_pptx

    long_text = "A" * 2500
    ct = SlidesContentTree(
        deck_title="Overflow Test",
        slides=[
            Slide(
                id="s1", slide_type="content", title="Overflow",
                elements=[SlideElement(id="e1", type="body", content=long_text)],
                speaker_notes="Check overflow.",
            ),
        ],
    )
    output = tmp_path / "overflow.pptx"
    export_to_pptx(ct, get_theme(), output)

    result = validate_pptx(output)
    assert result["layout_valid"] is False
    assert len(result["layout_errors"]) > 0


# === Phase 3: Quality Pass Acceptance Tests ===

def test_18_theme_catalog_reaches_100_plus() -> None:
    """Theme system provides 100+ total themes."""
    from core.studio.slides.themes import list_themes
    themes = list_themes(include_variants=True)
    assert len(themes) >= 100


def test_19_chart_slide_renders_native_chart(tmp_path) -> None:
    """Export a deck with bar chart spec produces PPTX with chart shape."""
    from pptx import Presentation as PptxPresentation
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.exporter import export_to_pptx
    from core.studio.slides.themes import get_theme

    ct = SlidesContentTree(
        deck_title="Chart Test",
        slides=[
            Slide(id="s1", slide_type="chart", title="Revenue",
                  elements=[SlideElement(id="e1", type="chart", content={
                      "chart_type": "bar",
                      "categories": ["Q1", "Q2", "Q3", "Q4"],
                      "series": [{"name": "Revenue", "values": [1.2, 1.8, 2.6, 3.1]}],
                  })],
                  speaker_notes="Discuss revenue trends."),
        ],
    )
    output = tmp_path / "chart_acceptance.pptx"
    export_to_pptx(ct, get_theme(), output)

    prs = PptxPresentation(str(output))
    chart_slide = prs.slides[0]
    has_chart = any(s.has_chart for s in chart_slide.shapes)
    assert has_chart, "Chart slide should contain native chart shape"


def test_20_notes_quality_passes_baseline() -> None:
    """repair_speaker_notes on sample deck achieves >= 90% pass rate."""
    from core.schemas.studio_schema import Slide, SlideElement, SlidesContentTree
    from core.studio.slides.notes import repair_speaker_notes, score_speaker_notes

    ct = SlidesContentTree(**_sample_content_tree_dict())
    repaired = repair_speaker_notes(ct)

    total = len(repaired.slides)
    pass_count = sum(
        1 for i, s in enumerate(repaired.slides)
        if score_speaker_notes(s, i, total)["passes"]
    )
    assert pass_count / total >= 0.90


def test_21_layout_quality_repaired_on_strict_export(tmp_path) -> None:
    """Export deck with huge body under strict mode repairs and succeeds."""
    import asyncio
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    art_id = str(uuid4())
    ct = _sample_content_tree_dict()
    # Inject overflow content — repair_layout() will truncate this
    ct["slides"][1]["elements"][0]["content"] = "X" * 2500
    artifact = Artifact(
        id=art_id, type=ArtifactType.slides, title="Overflow",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.pptx, strict_layout=True)
        )
    finally:
        loop.close()

    assert result["status"] == "completed"
    assert result["validator_results"] is not None


# --- Phase 4: Document DOCX/PDF export tests ---

def _sample_document_content_tree_dict() -> dict:
    """Return a sample document content tree dict for testing."""
    return {
        "doc_title": "Technical Specification: AI Agent Framework",
        "doc_type": "technical_spec",
        "abstract": "This document specifies the architecture and design of the AI Agent Framework.",
        "sections": [
            {
                "id": "sec1", "heading": "Introduction", "level": 1,
                "content": "The AI Agent Framework provides a modular architecture for building autonomous AI agents.",
                "subsections": [
                    {
                        "id": "sec1a", "heading": "Purpose", "level": 2,
                        "content": "This specification defines the core components and their interactions.",
                        "subsections": [], "citations": [],
                    }
                ],
                "citations": ["agent_paper"],
            },
            {
                "id": "sec2", "heading": "Architecture", "level": 1,
                "content": "The framework follows a layered architecture with clear separation of concerns.",
                "subsections": [], "citations": [],
            },
            {
                "id": "sec3", "heading": "Implementation", "level": 1,
                "content": "Implementation follows a phased approach with iterative delivery.",
                "subsections": [], "citations": [],
            },
        ],
        "bibliography": [
            {"key": "agent_paper", "title": "Autonomous AI Agents", "author": "Smith et al.", "year": "2024"},
        ],
        "metadata": {"audience": "engineers", "tone": "technical"},
    }


def test_22_docx_export_creates_valid_file(tmp_path):
    """Phase 4: DOCX export produces a valid, openable Word document."""
    from core.schemas.studio_schema import DocumentContentTree
    from core.studio.documents.exporter_docx import export_to_docx
    from core.studio.documents.validator import validate_docx

    ct_dict = _sample_document_content_tree_dict()
    ct = DocumentContentTree(**ct_dict)
    output_path = tmp_path / "test.docx"
    export_to_docx(ct, output_path)

    assert output_path.exists()
    assert output_path.stat().st_size > 0

    validation = validate_docx(output_path, ct)
    assert validation["valid"] is True
    assert validation["heading_count"] >= 3
    assert validation["text_present"] is True
    assert validation["bibliography_present"] is True


@pytest.mark.skipif(
    not _xhtml2pdf_available(),
    reason="xhtml2pdf not available",
)
def test_23_pdf_export_creates_valid_file(tmp_path):
    """Phase 4: PDF export produces a valid, openable PDF document."""
    from core.schemas.studio_schema import DocumentContentTree
    from core.studio.documents.exporter_pdf import export_to_pdf
    from core.studio.documents.validator import validate_pdf

    ct_dict = _sample_document_content_tree_dict()
    ct = DocumentContentTree(**ct_dict)
    output_path = tmp_path / "test.pdf"
    export_to_pdf(ct, output_path)

    assert output_path.exists()
    assert output_path.stat().st_size > 0

    validation = validate_pdf(output_path, ct)
    assert validation["valid"] is True
    assert validation["page_count"] >= 1
    assert validation["text_present"] is True


def test_24_document_export_job_lifecycle(tmp_path):
    """Phase 4: Full document export lifecycle — outline → draft → export DOCX."""
    import asyncio
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    art_id = str(uuid4())
    ct_dict = _sample_document_content_tree_dict()
    artifact = Artifact(
        id=art_id, type=ArtifactType.document, title="Tech Spec",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct_dict,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.docx)
        )
    finally:
        loop.close()

    assert result["status"] == "completed"
    assert result["format"] == "docx"
    assert result["output_uri"] is not None
    assert result["validator_results"]["valid"] is True


@pytest.mark.skipif(
    not _xhtml2pdf_available(),
    reason="xhtml2pdf not available",
)
def test_25_document_pdf_export_job_lifecycle(tmp_path):
    """Phase 4: Full document export lifecycle — export PDF."""
    import asyncio
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    art_id = str(uuid4())
    ct_dict = _sample_document_content_tree_dict()
    artifact = Artifact(
        id=art_id, type=ArtifactType.document, title="Tech Spec",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct_dict,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.pdf)
        )
    finally:
        loop.close()

    assert result["status"] == "completed"
    assert result["format"] == "pdf"
    assert result["validator_results"]["valid"] is True


# --- Phase 5: Sheet XLSX/CSV export tests ---

def _sample_sheet_content_tree_dict() -> dict:
    """Return a sample sheet content tree dict for testing."""
    return {
        "workbook_title": "Sales Report",
        "tabs": [
            {
                "id": "t1",
                "name": "Revenue",
                "headers": ["Month", "Amount", "Growth"],
                "rows": [
                    ["Jan", 1000, 0.0],
                    ["Feb", 1200, 0.2],
                    ["Mar", 1500, 0.25],
                    ["Apr", 1350, -0.1],
                ],
                "formulas": {"D2": "=C2/B2", "D3": "=C3/B3"},
                "column_widths": [100, 80, 80],
            },
            {
                "id": "t2",
                "name": "Summary",
                "headers": ["Metric", "Value"],
                "rows": [["Total", 5050], ["Average", 1262.5]],
                "formulas": {},
                "column_widths": [120, 80],
            },
        ],
    }


def test_26_xlsx_export_creates_valid_file(tmp_path) -> None:
    """Phase 5: XLSX export produces a valid, openable workbook."""
    from core.schemas.studio_schema import SheetContentTree
    from core.studio.sheets.exporter_xlsx import export_to_xlsx
    from core.studio.sheets.validator import validate_xlsx

    ct = SheetContentTree(**_sample_sheet_content_tree_dict())
    output_path = tmp_path / "test.xlsx"
    export_to_xlsx(ct, output_path)

    assert output_path.exists()
    assert output_path.stat().st_size > 0

    result = validate_xlsx(output_path, expected_sheet_names=["Revenue", "Summary"])
    assert result["valid"] is True
    assert result["sheet_count"] == 2


def test_27_csv_export_creates_valid_file(tmp_path) -> None:
    """Phase 5: CSV export produces a valid, openable CSV file."""
    from core.schemas.studio_schema import SheetContentTree
    from core.studio.sheets.exporter_csv import export_to_csv
    from core.studio.sheets.validator import validate_csv

    ct = SheetContentTree(**_sample_sheet_content_tree_dict())
    output_path = tmp_path / "test.csv"
    tab_name = export_to_csv(ct, output_path)

    assert output_path.exists()
    assert output_path.stat().st_size > 0
    assert tab_name == "Revenue"

    result = validate_csv(output_path, min_rows=4)
    assert result["valid"] is True
    assert result["column_count"] == 3


def test_28_upload_analysis_generates_summary_tabs(tmp_path) -> None:
    """Phase 5: CSV upload analysis generates analysis tabs and report."""
    import csv
    from core.studio.sheets.ingest import ingest_upload
    from core.studio.sheets.analysis import analyze_dataset, build_analysis_tabs

    csv_path = tmp_path / "data.csv"
    with open(csv_path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["Product", "Revenue", "Units"])
        for i in range(20):
            writer.writerow([f"Product_{i}", 100 + i * 10, 50 + i])

    dataset = ingest_upload("data.csv", csv_path.read_bytes(), "text/csv")
    assert dataset.columns == ["Product", "Revenue", "Units"]
    assert len(dataset.rows) == 20

    report = analyze_dataset(dataset)
    assert len(report.summary_stats) >= 2  # Revenue and Units

    tabs = build_analysis_tabs(dataset, report)
    tab_names = [t.name for t in tabs]
    assert "Uploaded_Data" in tab_names
    assert "Summary_Stats" in tab_names


def test_29_sheet_formula_refs_are_valid() -> None:
    """Phase 5: Formulas in generated sheet reference valid cells."""
    from core.schemas.studio_schema import SheetContentTree
    from core.studio.sheets.formulas import validate_tab_formulas

    ct = SheetContentTree(**_sample_sheet_content_tree_dict())
    for tab in ct.tabs:
        warnings = validate_tab_formulas(tab)
        # No formulas should be removed (all refs should be valid)
        for w in warnings:
            assert "removed" not in w.lower(), f"Valid formula was removed: {w}"


def test_30_sheet_export_job_lifecycle_completed(tmp_path) -> None:
    """Phase 5: Sheet export job transitions to completed with validator results."""
    import asyncio
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType, ExportFormat
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)

    art_id = str(uuid4())
    ct_dict = _sample_sheet_content_tree_dict()
    artifact = Artifact(
        id=art_id, type=ArtifactType.sheet, title="Sales Report",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct_dict,
    )
    storage.save_artifact(artifact)

    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.export_artifact(art_id, ExportFormat.xlsx)
        )
    finally:
        loop.close()

    assert result["status"] == "completed"
    assert result["format"] == "xlsx"
    assert result["output_uri"] is not None
    assert result["validator_results"]["valid"] is True
    assert result["validator_results"]["sheet_count"] >= 2


# === Phase 6: Edit Loop Acceptance Tests ===


def _make_slides_artifact(tmp_path):
    """Create a slides artifact with content tree for edit testing."""
    import asyncio
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage
    from core.studio.revision import RevisionManager

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)
    rm = RevisionManager(storage)

    art_id = str(uuid4())
    ct = _sample_content_tree_dict()
    artifact = Artifact(
        id=art_id,
        type=ArtifactType.slides,
        title="Edit Test",
        created_at=datetime.now(timezone.utc),
        updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    rev = rm.create_revision(art_id, ct, "Initial draft")
    artifact.revision_head_id = rev.id
    storage.save_artifact(artifact)
    return orch, storage, art_id


def _make_doc_artifact(tmp_path):
    """Create a document artifact with content tree for edit testing."""
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage
    from core.studio.revision import RevisionManager

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)
    rm = RevisionManager(storage)

    art_id = str(uuid4())
    ct = {
        "doc_title": "Test Report",
        "doc_type": "report",
        "abstract": "Summary.",
        "sections": [
            {"id": "sec1", "heading": "Introduction", "level": 1, "content": "Intro.", "subsections": [], "citations": []},
            {"id": "sec2", "heading": "Conclusion", "level": 1, "content": "End.", "subsections": [], "citations": []},
        ],
        "bibliography": [],
    }
    artifact = Artifact(
        id=art_id, type=ArtifactType.document, title="Edit Doc Test",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    rev = rm.create_revision(art_id, ct, "Initial draft")
    artifact.revision_head_id = rev.id
    storage.save_artifact(artifact)
    return orch, storage, art_id


def _make_sheet_artifact(tmp_path):
    """Create a sheet artifact with content tree for edit testing."""
    from datetime import datetime, timezone
    from uuid import uuid4
    from core.schemas.studio_schema import Artifact, ArtifactType
    from core.studio.orchestrator import ForgeOrchestrator
    from core.studio.storage import StudioStorage
    from core.studio.revision import RevisionManager

    storage = StudioStorage(base_dir=tmp_path / "studio")
    orch = ForgeOrchestrator(storage)
    rm = RevisionManager(storage)

    art_id = str(uuid4())
    ct = {
        "workbook_title": "Financial Model",
        "tabs": [
            {"id": "tab1", "name": "Revenue", "headers": ["Month", "MRR"],
             "rows": [["Jan", 5000], ["Feb", 5500]], "formulas": {}, "column_widths": [120, 100]},
        ],
    }
    artifact = Artifact(
        id=art_id, type=ArtifactType.sheet, title="Edit Sheet Test",
        created_at=datetime.now(timezone.utc), updated_at=datetime.now(timezone.utc),
        content_tree=ct,
    )
    rev = rm.create_revision(art_id, ct, "Initial draft")
    artifact.revision_head_id = rev.id
    storage.save_artifact(artifact)
    return orch, storage, art_id


def test_31_slides_edit_creates_valid_revision(tmp_path) -> None:
    """Editing a slides artifact creates a new revision."""
    import asyncio
    orch, storage, art_id = _make_slides_artifact(tmp_path)

    patch = {
        "artifact_type": "slides",
        "target": {"kind": "slide_index", "index": 2},
        "ops": [{"op": "SET", "path": "title", "value": "Edited Problem Statement"}],
        "summary": "Update slide 2 title",
    }
    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.edit_artifact(art_id, "Update slide 2 title", _patch_override=patch)
        )
    finally:
        loop.close()

    assert result["edit_result"]["status"] == "applied"
    revisions = storage.list_revisions(art_id)
    assert len(revisions) == 2  # initial + edit


def test_32_slides_edit_then_export_valid_pptx(tmp_path) -> None:
    """Edit then export produces a valid PPTX."""
    import asyncio
    from core.schemas.studio_schema import ExportFormat
    orch, storage, art_id = _make_slides_artifact(tmp_path)

    patch = {
        "artifact_type": "slides",
        "target": {"kind": "slide_index", "index": 3},
        "ops": [{"op": "SET", "path": "title", "value": "Revised Solution"}],
        "summary": "Update slide 3 title",
    }
    loop = asyncio.new_event_loop()
    try:
        loop.run_until_complete(orch.edit_artifact(art_id, "Edit", _patch_override=patch))
        export_result = loop.run_until_complete(orch.export_artifact(art_id, ExportFormat.pptx))
    finally:
        loop.close()

    assert export_result["status"] == "completed"


def test_33_document_edit_creates_valid_revision(tmp_path) -> None:
    """Editing a document artifact creates a new revision."""
    import asyncio
    orch, storage, art_id = _make_doc_artifact(tmp_path)

    patch = {
        "artifact_type": "document",
        "target": {"kind": "section_id", "id": "sec1"},
        "ops": [{"op": "SET", "path": "content", "value": "Updated introduction content."}],
        "summary": "Update intro section",
    }
    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.edit_artifact(art_id, "Update intro", _patch_override=patch)
        )
    finally:
        loop.close()

    assert result["edit_result"]["status"] == "applied"


def test_34_sheet_edit_creates_valid_revision(tmp_path) -> None:
    """Editing a sheet artifact creates a new revision."""
    import asyncio
    orch, storage, art_id = _make_sheet_artifact(tmp_path)

    patch = {
        "artifact_type": "sheet",
        "target": {"kind": "tab_name", "name": "Revenue"},
        "ops": [{"op": "SET", "path": "rows[0][1]", "value": 9999}],
        "summary": "Update revenue value",
    }
    loop = asyncio.new_event_loop()
    try:
        result = loop.run_until_complete(
            orch.edit_artifact(art_id, "Update revenue", _patch_override=patch)
        )
    finally:
        loop.close()

    assert result["edit_result"]["status"] == "applied"


def test_35_edit_then_export_all_types(tmp_path) -> None:
    """Edit + export works for slides, documents, and sheets."""
    import asyncio
    from core.schemas.studio_schema import ExportFormat

    # Slides
    s_orch, _, s_id = _make_slides_artifact(tmp_path)
    # Document
    d_orch, _, d_id = _make_doc_artifact(tmp_path)
    # Sheet
    sh_orch, _, sh_id = _make_sheet_artifact(tmp_path)

    loop = asyncio.new_event_loop()
    try:
        # Edit all
        loop.run_until_complete(s_orch.edit_artifact(
            s_id, "edit", _patch_override={
                "artifact_type": "slides",
                "target": {"kind": "slide_index", "index": 1},
                "ops": [{"op": "SET", "path": "title", "value": "New Title"}],
                "summary": "edit",
            }
        ))
        loop.run_until_complete(d_orch.edit_artifact(
            d_id, "edit", _patch_override={
                "artifact_type": "document",
                "target": {"kind": "section_id", "id": "sec1"},
                "ops": [{"op": "SET", "path": "content", "value": "New content."}],
                "summary": "edit",
            }
        ))
        loop.run_until_complete(sh_orch.edit_artifact(
            sh_id, "edit", _patch_override={
                "artifact_type": "sheet",
                "target": {"kind": "tab_name", "name": "Revenue"},
                "ops": [{"op": "SET", "path": "headers", "value": ["Month", "MRR", "Growth"]}],
                "summary": "edit",
            }
        ))

        # Export all
        s_job = loop.run_until_complete(s_orch.export_artifact(s_id, ExportFormat.pptx))
        d_job = loop.run_until_complete(d_orch.export_artifact(d_id, ExportFormat.docx))
        sh_job = loop.run_until_complete(sh_orch.export_artifact(sh_id, ExportFormat.xlsx))
    finally:
        loop.close()

    assert s_job["status"] == "completed"
    assert d_job["status"] == "completed"
    assert sh_job["status"] == "completed"
