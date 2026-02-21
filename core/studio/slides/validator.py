"""Open-validation for exported PPTX files — v2 with quality checks."""

from pathlib import Path
from typing import Any, Dict, Optional


def validate_pptx(
    file_path: Path,
    expected_slide_count: int | None = None,
    content_tree: Optional[Any] = None,
) -> Dict[str, Any]:
    """Validate a PPTX file with structural, layout, chart, notes, and content checks.

    Returns a dict with valid, slide_count, has_notes, errors,
    layout_valid, layout_warnings, layout_errors,
    notes_quality_valid, chart_quality_valid, quality_score.
    """
    errors = []
    slide_count = 0
    has_notes = False
    layout_warnings = []
    layout_errors = []
    notes_quality_valid = True
    chart_quality_valid = True

    # === Layer 1: Structural ===
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slide_count = len(prs.slides)

        if expected_slide_count is not None and slide_count != expected_slide_count:
            errors.append(
                f"Slide count mismatch: expected {expected_slide_count}, got {slide_count}"
            )

        for slide in prs.slides:
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame.text.strip():
                    has_notes = True
                    break
            except Exception:
                continue

    except Exception as e:
        errors.append(f"Failed to open PPTX: {str(e)}")
        return _build_result(
            errors=errors, slide_count=slide_count, has_notes=has_notes,
            layout_warnings=layout_warnings, layout_errors=layout_errors,
            notes_quality_valid=notes_quality_valid,
            chart_quality_valid=chart_quality_valid,
        )

    # === Layer 2: Layout (blocking) ===
    from core.studio.slides.exporter import SLIDE_WIDTH, SLIDE_HEIGHT
    BLOCK_CHAR_LIMIT = 800
    SLIDE_CHAR_LIMIT = 1600

    try:
        from pptx import Presentation as _Prs
        _prs = _Prs(str(file_path))
        for slide_idx, slide in enumerate(_prs.slides):
            slide_total = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_len = len(shape.text_frame.text)
                    slide_total += text_len
                    if text_len > BLOCK_CHAR_LIMIT:
                        layout_errors.append(
                            f"Slide {slide_idx + 1}: text block exceeds {BLOCK_CHAR_LIMIT} chars ({text_len} chars)"
                        )
                # Out-of-bounds detection
                try:
                    if hasattr(shape, 'left') and hasattr(shape, 'width'):
                        if shape.left + shape.width > SLIDE_WIDTH:
                            layout_errors.append(
                                f"Slide {slide_idx + 1}: shape extends beyond slide width"
                            )
                    if hasattr(shape, 'top') and hasattr(shape, 'height'):
                        if shape.top + shape.height > SLIDE_HEIGHT:
                            layout_errors.append(
                                f"Slide {slide_idx + 1}: shape extends beyond slide height"
                            )
                except Exception:
                    pass

                # Small font check (advisory)
                if shape.has_text_frame:
                    try:
                        for p in shape.text_frame.paragraphs:
                            if p.font and p.font.size and p.font.size < 127000:  # < 10pt in EMUs
                                layout_warnings.append(
                                    f"Slide {slide_idx + 1}: text with font size < 10pt"
                                )
                                break
                    except Exception:
                        pass

            if slide_total > SLIDE_CHAR_LIMIT:
                layout_errors.append(
                    f"Slide {slide_idx + 1}: total text density exceeds {SLIDE_CHAR_LIMIT} chars ({slide_total} chars)"
                )
    except Exception:
        pass

    # === Layer 3: Chart quality (advisory) ===
    if content_tree is not None:
        try:
            from pptx import Presentation as _Prs3
            _prs3 = _Prs3(str(file_path))
            ct_slides = content_tree.slides if hasattr(content_tree, 'slides') else []

            for slide_idx, ct_slide in enumerate(ct_slides):
                if ct_slide.slide_type == "chart" and slide_idx < len(_prs3.slides):
                    pptx_slide = _prs3.slides[slide_idx]
                    has_chart_shape = False
                    has_chart_marker = False

                    for shape in pptx_slide.shapes:
                        if shape.has_chart:
                            has_chart_shape = True
                            break
                        if shape.has_text_frame and "[Chart:" in shape.text_frame.text:
                            has_chart_marker = True

                    if not has_chart_shape and not has_chart_marker:
                        chart_quality_valid = False
        except Exception:
            pass

    # === Layer 4: Notes quality (advisory) ===
    if content_tree is not None:
        try:
            from core.studio.slides.notes import score_speaker_notes
            ct_slides = content_tree.slides if hasattr(content_tree, 'slides') else []
            total = len(ct_slides)
            if total > 0:
                pass_count = 0
                has_empty = False
                for i, ct_slide in enumerate(ct_slides):
                    score = score_speaker_notes(ct_slide, i, total)
                    if score["passes"]:
                        pass_count += 1
                    if score["is_empty"]:
                        has_empty = True

                pass_rate = pass_count / total
                if pass_rate < 0.90 or has_empty:
                    notes_quality_valid = False
        except Exception:
            pass

    # === Layer 5: Content heuristics (advisory) ===
    if content_tree is not None:
        try:
            ct_slides = content_tree.slides if hasattr(content_tree, 'slides') else []
            for slide_idx, ct_slide in enumerate(ct_slides):
                # Title length check
                if ct_slide.title and len(ct_slide.title) > 60:
                    layout_warnings.append(
                        f"Slide {slide_idx + 1}: title exceeds 60 chars ({len(ct_slide.title)} chars)"
                    )

                # Bullet count check
                for el in ct_slide.elements:
                    if el.type == "bullet_list" and isinstance(el.content, list) and len(el.content) > 7:
                        layout_warnings.append(
                            f"Slide {slide_idx + 1}: bullet list has {len(el.content)} items"
                        )

                # Sparse content check (non-title slides)
                if ct_slide.slide_type != "title":
                    total_text = 0
                    for el in ct_slide.elements:
                        if isinstance(el.content, str):
                            total_text += len(el.content)
                        elif isinstance(el.content, list):
                            total_text += sum(len(str(item)) for item in el.content)
                    if total_text < 20:
                        layout_warnings.append(
                            f"Slide {slide_idx + 1}: sparse content ({total_text} chars)"
                        )
        except Exception:
            pass

    return _build_result(
        errors=errors, slide_count=slide_count, has_notes=has_notes,
        layout_warnings=layout_warnings, layout_errors=layout_errors,
        notes_quality_valid=notes_quality_valid,
        chart_quality_valid=chart_quality_valid,
    )


def _build_result(
    *,
    errors: list,
    slide_count: int,
    has_notes: bool,
    layout_warnings: list,
    layout_errors: list,
    notes_quality_valid: bool,
    chart_quality_valid: bool,
) -> Dict[str, Any]:
    """Build the standardized validator result dict."""
    layout_valid = len(layout_errors) == 0

    # Quality score calculation
    score = 100
    if not layout_valid:
        score -= 40
    score -= min(len(layout_warnings), 4) * 5
    if not notes_quality_valid:
        score -= 20
    if not chart_quality_valid:
        score -= 20
    quality_score = max(0, min(100, score))

    return {
        "valid": len(errors) == 0,
        "slide_count": slide_count,
        "has_notes": has_notes,
        "errors": errors,
        "layout_valid": layout_valid,
        "layout_warnings": layout_warnings,
        "layout_errors": layout_errors,
        "notes_quality_valid": notes_quality_valid,
        "chart_quality_valid": chart_quality_valid,
        "quality_score": quality_score,
    }
