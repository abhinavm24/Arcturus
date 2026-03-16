"""PPTX renderer for Forge slides — programmatic shapes, no templates."""

import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import hashlib
import logging

from core.schemas.studio_schema import SlideTheme, SlidesContentTree
from core.studio.slides.themes import _blend_color
from core.studio.slides.images import extract_image_url

_logger = logging.getLogger(__name__)


def _download_image_url(url: str) -> io.BytesIO | None:
    """Download an image from a URL and return as BytesIO, or None on failure."""
    try:
        import urllib.request
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = resp.read()
            if len(data) > 0:
                return io.BytesIO(data)
    except Exception as e:
        _logger.warning("Failed to download image from %s: %s", url, e)
    return None

# 16:9 widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margins
MARGIN_LEFT = Inches(0.75)
MARGIN_TOP = Inches(0.75)
MARGIN_RIGHT = Inches(0.75)
MARGIN_BOTTOM = Inches(0.5)

# Content area
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM

# Title area
TITLE_TOP = Inches(0.5)
TITLE_LEFT = MARGIN_LEFT
TITLE_WIDTH = CONTENT_WIDTH
TITLE_HEIGHT = Inches(1.0)

# Body area (below title)
BODY_TOP = Inches(1.8)
BODY_LEFT = MARGIN_LEFT
BODY_WIDTH = CONTENT_WIDTH
BODY_HEIGHT = Inches(5.0)

# Two-column split
COLUMN_GAP = Inches(0.5)
COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2

# Chart area (used when slide_type == "chart")
CHART_TOP = Inches(2.0)
CHART_LEFT = MARGIN_LEFT
CHART_WIDTH = CONTENT_WIDTH
CHART_HEIGHT = Inches(3.8)

# Caption below chart
CAPTION_TOP = Inches(6.0)
CAPTION_HEIGHT = Inches(1.0)

# === Design Tokens ===

_DESIGN_TOKENS = {
    # Typography scale
    "title_size": Pt(44),
    "heading_size": Pt(32),
    "subheading_size": Pt(24),
    "body_size": Pt(18),
    "body_small_size": Pt(16),
    "code_size": Pt(14),
    "footer_size": Pt(10),
    "quote_size": Pt(28),
    "attribution_size": Pt(16),
    "stat_callout_size": Pt(56),
    "stat_label_size": Pt(14),

    # Text frame margins
    "margin_left": Inches(0.18),
    "margin_right": Inches(0.18),
    "margin_top": Inches(0.12),
    "margin_bottom": Inches(0.12),

    # Paragraph spacing
    "line_spacing": 1.40,
    "para_space_after": Pt(10),
    "bullet_space_after": Pt(12),

    # Bullet formatting
    "bullet_indent": Inches(0.25),
    "bullet_hanging": Inches(0.20),

    # Kicker / Takeaway
    "kicker_size": Pt(14),
    "kicker_char_spacing": Pt(1.2),
    "takeaway_size": Pt(16),

    # Caption tokens
    "caption_size": Pt(14),

    # Slide chrome
    "accent_bar_height": Inches(0.10),
    "accent_bar_top": Inches(7.1),
    "footer_height": Inches(0.25),
    "footer_top": Inches(7.2),
}

# Font scale compensation for serif fonts
_FONT_SCALE = {
    "Garamond": 1.20,
    "Book Antiqua": 1.10,
    "Georgia": 1.04,
    "Constantia": 1.08,
    "Cambria": 1.06,
}

# Chart type mapping
try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "funnel": XL_CHART_TYPE.BAR_CLUSTERED,
    }
    _CHARTS_AVAILABLE = True
except ImportError:
    _CHARTS_AVAILABLE = False
    _CHART_TYPE_MAP = {}


def export_to_pptx(
    content_tree: SlidesContentTree,
    theme: SlideTheme,
    output_path: Path,
    images: dict[str, io.BytesIO] | None = None,
) -> Path:
    """Export a SlidesContentTree to PPTX format.

    Args:
        images: Optional dict mapping slide ID to JPEG BytesIO buffers for
                image_text slides. When provided, actual images are embedded
                instead of text placeholders.

    Returns the output file path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout
    total_slides = len(content_tree.slides)
    section_count = 0

    for i, slide_data in enumerate(content_tree.slides):
        pptx_slide = prs.slides.add_slide(blank_layout)

        if slide_data.slide_type == "section_divider":
            section_count += 1

        renderer = _RENDERERS.get(slide_data.slide_type, _render_content)
        renderer(pptx_slide, slide_data, theme, images=images, slide_index=i,
                 section_number=section_count, total_slides=total_slides)

        # Add slide chrome (skip first and last slides)
        if 0 < i < total_slides - 1:
            _add_slide_chrome(pptx_slide, theme, i + 1, total_slides)

        if slide_data.speaker_notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# === Helper Functions ===

def _scaled_font_size(base_size, font_name):
    """Apply font scale compensation for serif fonts."""
    scale = _FONT_SCALE.get(font_name, 1.0)
    if scale == 1.0:
        return base_size
    return Pt(int(base_size.pt * scale + 0.5))


def _compute_body_top(title: str, font_size_pt: float = 32.0,
                      has_kicker: bool = False) -> tuple:
    """Compute adjusted body top and title font size based on title length.

    Returns (body_top_inches, title_font_size_pt, title_top_inches).
    Long titles that wrap to multiple lines push the body area down.
    When a kicker is present, the title shifts down to avoid overlap.
    """
    title = (title or "").strip()
    title_len = len(title)

    # Kicker occupies ~0.35" at the top; shift title down when present
    kicker_offset = 0.35 if has_kicker else 0.0

    # Approximate characters per line at given font size on CONTENT_WIDTH (~11.83")
    # At 36pt ~2 chars/inch → ~24 chars/line; at 28pt ~2.5 chars/inch → ~30 chars/line
    if title_len > 80:
        # 3+ lines likely — shrink title font and push body down
        adjusted_font = 28.0
        body_top = 2.3 + kicker_offset
    elif title_len > 50:
        # 2 lines likely — push body down slightly
        adjusted_font = font_size_pt
        body_top = 2.1 + kicker_offset
    else:
        # Single line — default
        adjusted_font = font_size_pt
        body_top = 1.8 + kicker_offset

    # Ensure body area stays >= 4.2" (slide bottom at ~7.1" minus body_top)
    max_body_top = 7.1 - 4.2  # = 2.9"
    body_top = min(body_top, max_body_top)

    title_top = TITLE_TOP + Inches(kicker_offset)

    return Inches(body_top), Pt(int(adjusted_font)), title_top


_MD_INLINE_RE = re.compile(
    r'\*\*\*(?!\s)(.+?)(?<!\s)\*\*\*'   # ***bold+italic*** (no inner spaces at delimiters)
    r'|\*\*(?!\s)(.+?)(?<!\s)\*\*'       # **bold**
    r'|\*(?!\s)(.+?)(?<!\s)\*'           # *italic*
    r'|==(?!\s)(.+?)(?<!\s)=='           # ==highlight==
)


def _parse_markdown_runs(text: str) -> list[tuple[str, bool, bool, bool]]:
    """Parse markdown inline formatting into (text, bold, italic, highlight) segments."""
    if not text:
        return [("", False, False, False)]
    segments: list[tuple[str, bool, bool, bool]] = []
    last_end = 0
    for m in _MD_INLINE_RE.finditer(text):
        # Plain text before this match
        if m.start() > last_end:
            segments.append((text[last_end:m.start()], False, False, False))
        if m.group(1) is not None:       # ***bold+italic***
            segments.append((m.group(1), True, True, False))
        elif m.group(2) is not None:     # **bold**
            segments.append((m.group(2), True, False, False))
        elif m.group(3) is not None:     # *italic*
            segments.append((m.group(3), False, True, False))
        elif m.group(4) is not None:     # ==highlight==
            segments.append((m.group(4), False, False, True))
        last_end = m.end()
    # Trailing plain text (or entire string if no matches)
    if last_end < len(text):
        segments.append((text[last_end:], False, False, False))
    return segments or [("", False, False, False)]


def _apply_markdown_runs(paragraph, text, *, font_name, font_size, font_color,
                         bold=False, accent_color=None):
    """Replace paragraph text with per-run markdown formatting."""
    if isinstance(font_color, str):
        font_color = RGBColor.from_string(font_color.lstrip("#"))
    if isinstance(accent_color, str):
        accent_color = RGBColor.from_string(accent_color.lstrip("#"))

    segments = _parse_markdown_runs(str(text))

    # Clear existing text
    paragraph.clear()

    # Set paragraph-level font size for validator compatibility
    paragraph.font.size = font_size

    for seg_text, seg_bold, seg_italic, seg_highlight in segments:
        run = paragraph.add_run()
        run.text = seg_text
        run.font.name = font_name
        run.font.size = font_size
        if seg_highlight and accent_color is not None:
            run.font.color.rgb = accent_color
            run.font.bold = True
        else:
            run.font.color.rgb = font_color
            run.font.bold = bold or seg_bold
        run.font.italic = seg_italic


def _add_text_box(slide, text, left, top, width, height,
                  font_name="Calibri", font_size=None,
                  font_color="#000000", alignment=PP_ALIGN.LEFT,
                  bold=False, parse_markdown=True, accent_color=None):
    """Add a text box shape with styled text and design token margins."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Apply margins from design tokens
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    p = tf.paragraphs[0]
    p.alignment = alignment
    p.line_spacing = _DESIGN_TOKENS["line_spacing"]
    p.space_after = _DESIGN_TOKENS["para_space_after"]

    if parse_markdown:
        _apply_markdown_runs(
            p, text, font_name=font_name, font_size=scaled_size,
            font_color=font_color, bold=bold, accent_color=accent_color,
        )
    else:
        p.text = str(text)
        p.font.name = font_name
        p.font.size = scaled_size
        p.font.bold = bold
        if isinstance(font_color, str):
            font_color = RGBColor.from_string(font_color.lstrip("#"))
        p.font.color.rgb = font_color


def _add_bullet_list(slide, items, left, top, width, height,
                     font_name="Calibri", font_size=None,
                     font_color="#000000", accent_color=None):
    """Add a text box with bullet-point paragraphs."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_small_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    if isinstance(font_color, str):
        resolved_color = RGBColor.from_string(font_color.lstrip("#"))
    else:
        resolved_color = font_color

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Apply margins
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = _DESIGN_TOKENS["bullet_space_after"]
        p.line_spacing = _DESIGN_TOKENS["line_spacing"]
        _apply_markdown_runs(
            p, f"\u2022 {item}", font_name=font_name,
            font_size=scaled_size, font_color=resolved_color,
            accent_color=accent_color,
        )


def _find_element(slide_data, element_type):
    """Find first element of a given type in slide data."""
    for el in slide_data.elements:
        if el.type == element_type:
            return el
    return None


def _has_kicker(slide_data):
    """Check if slide data contains a non-empty kicker element."""
    el = _find_element(slide_data, "kicker")
    return el is not None and bool(el.content)


def _render_kicker(slide, slide_data, theme):
    """Render a kicker element above the title (small uppercase text in accent color).

    Returns True if a kicker was rendered, False otherwise.
    """
    kicker_el = _find_element(slide_data, "kicker")
    if not kicker_el or not kicker_el.content:
        return False
    kicker_text = str(kicker_el.content).upper()
    _add_text_box(slide, kicker_text,
                  left=TITLE_LEFT, top=Inches(0.25),
                  width=TITLE_WIDTH, height=Inches(0.3),
                  font_name=theme.font_body,
                  font_size=_DESIGN_TOKENS["kicker_size"],
                  font_color=theme.colors.accent,
                  bold=True, parse_markdown=False)
    # Apply letter-spacing to kicker text runs
    if slide.shapes:
        last_shape = slide.shapes[-1]
        if last_shape.has_text_frame:
            spc_val = str(int(_DESIGN_TOKENS["kicker_char_spacing"].pt * 100))
            for p in last_shape.text_frame.paragraphs:
                for run in p.runs:
                    rPr = run._r.get_or_add_rPr()
                    rPr.set("spc", spc_val)
    return True


_TAKEAWAY_HEIGHT = Inches(0.6)
_TAKEAWAY_GAP = Inches(0.1)  # Gap between body content and takeaway
_TAKEAWAY_RESERVE = _TAKEAWAY_HEIGHT + _TAKEAWAY_GAP  # Total space to reserve


def _has_takeaway(slide_data):
    """Check if slide data contains a non-empty takeaway element."""
    el = _find_element(slide_data, "takeaway")
    return el is not None and bool(el.content)


def _render_takeaway(slide, slide_data, theme, top=None):
    """Render a takeaway element at the bottom of the slide.

    Returns True if a takeaway was rendered, False otherwise.
    """
    takeaway_el = _find_element(slide_data, "takeaway")
    if not takeaway_el or not takeaway_el.content:
        return False
    takeaway_top = top or Inches(6.4)
    _add_text_box(slide, str(takeaway_el.content),
                  left=TITLE_LEFT, top=takeaway_top,
                  width=TITLE_WIDTH, height=_TAKEAWAY_HEIGHT,
                  font_name=theme.font_body,
                  font_size=_DESIGN_TOKENS["takeaway_size"],
                  font_color=theme.colors.secondary,
                  bold=True)
    return True


def _css_color_to_hex(css_color: str) -> str | None:
    """Convert a CSS color (hex, rgb, rgba) to a 6-digit hex string."""
    if not css_color:
        return None
    css_color = css_color.strip()
    # 6-digit hex
    m = re.match(r'^#([0-9a-fA-F]{6})$', css_color)
    if m:
        return css_color.upper()
    # 3-digit hex
    m = re.match(r'^#([0-9a-fA-F]{3})$', css_color)
    if m:
        c = m.group(1)
        return f"#{c[0]*2}{c[1]*2}{c[2]*2}".upper()
    # 8-digit hex (with alpha) — take first 6
    m = re.match(r'^#([0-9a-fA-F]{6})[0-9a-fA-F]{2}$', css_color)
    if m:
        return f"#{m.group(1)}".upper()
    # rgb(r, g, b) or rgba(r, g, b, a)
    m = re.match(r'rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', css_color)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        return f"#{r:02X}{g:02X}{b:02X}"
    return None


_LINEAR_GRADIENT_RE = re.compile(
    r'linear-gradient\s*\(\s*'
    r'(?:(\d+)deg\s*,\s*)?'
    r'(?:to\s+\w+(?:\s+\w+)?\s*,\s*)?'
    r'(#[0-9a-fA-F]{3,8}|rgba?\([^)]+\))'
    r'(?:\s+\d+%)?\s*,\s*'
    r'(?:[^,]+,\s*)*?'  # skip middle stops
    r'(#[0-9a-fA-F]{3,8}|rgba?\([^)]+\))'
    r'(?:\s+\d+%)?'
    r'\s*\)',
    re.IGNORECASE,
)


def _parse_linear_gradient(css_value: str) -> tuple[str, str, int] | None:
    """Parse CSS linear-gradient into (color1_hex, color2_hex, angle).

    Uses first and last color stops (PPTX only supports 2 stops).
    Returns None if not parseable.
    """
    m = _LINEAR_GRADIENT_RE.search(css_value)
    if not m:
        return None
    angle = int(m.group(1)) if m.group(1) else 135
    c1 = _css_color_to_hex(m.group(2))
    c2 = _css_color_to_hex(m.group(3))
    if c1 and c2:
        # Convert CSS angle (0=up, clockwise) to PPTX angle
        pptx_angle = (450 - angle) % 360
        return (c1, c2, pptx_angle)
    return None


def _resolve_title_color(slide_data, theme) -> str:
    """Get title color from slide_style override or theme."""
    meta = getattr(slide_data, "metadata", None) or {}
    ss = meta.get("slide_style") or meta.get("visual_style") or {}
    if isinstance(ss, dict):
        title_style = ss.get("title", {})
        if isinstance(title_style, dict):
            color = title_style.get("color")
            if color:
                hex_color = _css_color_to_hex(color)
                if hex_color:
                    return hex_color
    return theme.colors.primary


def _resolve_body_color(slide_data, theme) -> str:
    """Get body text color from slide_style override or theme."""
    meta = getattr(slide_data, "metadata", None) or {}
    ss = meta.get("slide_style") or meta.get("visual_style") or {}
    if isinstance(ss, dict):
        body_style = ss.get("body", {})
        if isinstance(body_style, dict):
            color = body_style.get("color")
            if color:
                hex_color = _css_color_to_hex(color)
                if hex_color:
                    return hex_color
    return theme.colors.text


def _set_slide_background(slide, theme, slide_style=None):
    """Set slide background from slide_style CSS or theme default.

    Supports:
    - LLM-specified CSS backgrounds (solid hex, linear-gradient)
    - Legacy visual_style bg_variant values
    - Theme background fallback
    """
    bg_value = None
    if isinstance(slide_style, dict):
        # New format: slide_style.background.value
        bg_obj = slide_style.get("background")
        if isinstance(bg_obj, dict):
            bg_value = bg_obj.get("value")
        # Legacy format: visual_style.bg_variant — convert to CSS
        if not bg_value:
            bg_variant = slide_style.get("bg_variant")
            if bg_variant == "dark_invert":
                title_bg = getattr(theme.colors, "title_background", None) or _blend_color(theme.colors.background, "#000000", 0.85)
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(title_bg.lstrip("#"))
                return
            elif bg_variant == "accent_wash":
                washed = _blend_color(theme.colors.background, theme.colors.accent, 0.05)
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(washed.lstrip("#"))
                return
            elif bg_variant == "gradient":
                bg_value = f"linear-gradient(135deg, {theme.colors.background}, {_blend_color(theme.colors.background, theme.colors.primary, 0.08)})"

    bg_hex = theme.colors.background

    if bg_value:
        # Try to parse as gradient
        gradient = _parse_linear_gradient(bg_value)
        if gradient:
            c1, c2, angle = gradient
            try:
                fill = slide.background.fill
                fill.gradient()
                fill.gradient_stops[0].color.rgb = RGBColor.from_string(c1.lstrip("#"))
                fill.gradient_stops[1].color.rgb = RGBColor.from_string(c2.lstrip("#"))
                fill.gradient_angle = angle
                return
            except Exception:
                pass  # Fall through to solid

        # Try to extract a solid hex color
        hex_color = _css_color_to_hex(bg_value)
        if not hex_color:
            # Try first hex in the string
            hex_match = re.search(r'#([0-9a-fA-F]{6})', bg_value)
            if hex_match:
                hex_color = f"#{hex_match.group(1)}"
        if hex_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))
            return

    # Theme-level background style
    if getattr(theme, "background_style", None) == "gradient":
        try:
            blended = _blend_color(bg_hex, theme.colors.primary, 0.08)
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))
            fill.gradient_stops[1].color.rgb = RGBColor.from_string(blended.lstrip("#"))
            fill.gradient_angle = 270
            return
        except Exception:
            pass

    # Fallback: solid theme background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(bg_hex.lstrip("#"))


# === Slide Chrome ===

def _add_slide_chrome(slide, theme, slide_number, total_slides):
    """Add progress bar and slide number to a slide."""
    bar_height = Inches(0.04)
    bar_top = _DESIGN_TOKENS["accent_bar_top"]

    # Background bar — full width, light tint
    bg_tint = _blend_color(theme.colors.text_light, theme.colors.background, 0.20)
    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), bar_top, SLIDE_WIDTH, bar_height,
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor.from_string(bg_tint.lstrip("#"))
    bg_bar.line.fill.background()

    # Fill bar — proportional width, accent color
    progress = slide_number / total_slides
    fill_width = int(SLIDE_WIDTH * progress)
    fill_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), bar_top, fill_width, bar_height,
    )
    fill_bar.fill.solid()
    fill_bar.fill.fore_color.rgb = RGBColor.from_string(
        theme.colors.accent.lstrip("#")
    )
    fill_bar.line.fill.background()

    # Slide number — right-aligned (just the number)
    num_box = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN_RIGHT - Inches(1.2), _DESIGN_TOKENS["footer_top"],
        Inches(1.2), _DESIGN_TOKENS["footer_height"],
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = _DESIGN_TOKENS["footer_size"]
    p.font.color.rgb = RGBColor.from_string(
        theme.colors.text_light.lstrip("#")
    )


# === Icon Bullets ===

_BULLET_CHAR = "\u25CF"  # ● — solid circle bullet


def _add_icon_bullet_list(slide, items, left, top, width, height,
                          font_name="Calibri", font_size=None,
                          font_color="#000000", icon_color="#4472C4",
                          accent_color=None):
    """Add a bullet list with inline colored bullet characters."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["body_size"]
    scaled_size = _scaled_font_size(font_size, font_name)
    if isinstance(font_color, str):
        resolved_color = RGBColor.from_string(font_color.lstrip("#"))
    else:
        resolved_color = font_color
    if isinstance(icon_color, str):
        resolved_icon_color = RGBColor.from_string(icon_color.lstrip("#"))
    else:
        resolved_icon_color = icon_color

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = _DESIGN_TOKENS["margin_left"]
    tf.margin_right = _DESIGN_TOKENS["margin_right"]
    tf.margin_top = _DESIGN_TOKENS["margin_top"]
    tf.margin_bottom = _DESIGN_TOKENS["margin_bottom"]

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = _DESIGN_TOKENS["bullet_space_after"]
        p.line_spacing = _DESIGN_TOKENS["line_spacing"]

        # Clear default text, then add bullet + markdown runs manually
        # (cannot use _apply_markdown_runs because it calls paragraph.clear()
        # which would wipe the bullet character)
        p.clear()
        p.font.size = scaled_size

        # Colored bullet character as first run
        bullet_run = p.add_run()
        bullet_run.text = _BULLET_CHAR + "  "
        bullet_run.font.size = scaled_size
        bullet_run.font.name = font_name
        bullet_run.font.color.rgb = resolved_icon_color

        # Item text with per-run markdown formatting
        if isinstance(accent_color, str):
            resolved_accent = RGBColor.from_string(accent_color.lstrip("#"))
        else:
            resolved_accent = accent_color
        for seg_text, seg_bold, seg_italic, seg_highlight in _parse_markdown_runs(str(item)):
            run = p.add_run()
            run.text = seg_text
            run.font.name = font_name
            run.font.size = scaled_size
            if seg_highlight and resolved_accent is not None:
                run.font.color.rgb = resolved_accent
                run.font.bold = True
            else:
                run.font.color.rgb = resolved_color
                run.font.bold = seg_bold
            run.font.italic = seg_italic


# === Decorative Accents ===

def _add_decorative_accent(slide, theme, slide_index):
    """Add a small decorative visual accent to text-heavy slides.

    Rotates through 3 styles based on slide_index to add visual variety.
    """
    accent_hex = theme.colors.accent.lstrip("#")
    style = slide_index % 3

    if style == 0:
        # Accent dot cluster — 3 small circles in top-right corner
        for i in range(3):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                SLIDE_WIDTH - Inches(0.9 + i * 0.25), Inches(0.3),
                Inches(0.15), Inches(0.15),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
            dot.line.fill.background()
    elif style == 1:
        # Corner bracket — thin L-shape lines in bottom-left
        h_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), SLIDE_HEIGHT - Inches(0.7),
            Inches(1.2), Inches(0.04),
        )
        h_line.fill.solid()
        h_line.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        h_line.line.fill.background()
        v_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), SLIDE_HEIGHT - Inches(1.7),
            Inches(0.04), Inches(1.04),
        )
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        v_line.line.fill.background()
    else:
        # Horizontal divider rule below title
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, Inches(1.65), BODY_WIDTH, Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string(accent_hex)
        rule.line.fill.background()


# === Background Geometric Graphics ===

# Shape configs for title, closing, and section divider slides.
# Each entry: shape type, position (left, top, width, height in Inches),
# color_key (theme attribute), alpha_pct (0–100 fill opacity).
_BG_GRAPHICS = {
    "title": [
        # Large circle bottom-right — anchors the corner
        {"shape": "circle", "left": 10.2, "top": 4.5, "w": 4.0, "h": 4.0,
         "color_key": "primary", "alpha_pct": 20},
        # Medium circle top-left — balances the composition
        {"shape": "circle", "left": -1.0, "top": -0.8, "w": 2.8, "h": 2.8,
         "color_key": "accent", "alpha_pct": 18},
        # Small accent dot upper-right
        {"shape": "circle", "left": 11.8, "top": 0.4, "w": 1.0, "h": 1.0,
         "color_key": "accent", "alpha_pct": 25},
        # Ring bottom-left — open shape for contrast
        {"shape": "ring", "left": 0.2, "top": 5.2, "w": 2.2, "h": 2.2,
         "color_key": "accent", "alpha_pct": 15},
    ],
    "closing": [
        # Large circle top-right (mirrored from title)
        {"shape": "circle", "left": 9.8, "top": -1.5, "w": 4.0, "h": 4.0,
         "color_key": "primary", "alpha_pct": 20},
        # Medium circle bottom-left
        {"shape": "circle", "left": -0.8, "top": 5.0, "w": 2.8, "h": 2.8,
         "color_key": "accent", "alpha_pct": 18},
        # Small accent dot lower-left
        {"shape": "circle", "left": 1.2, "top": 6.0, "w": 1.0, "h": 1.0,
         "color_key": "accent", "alpha_pct": 25},
        # Ring top-left
        {"shape": "ring", "left": -0.2, "top": -0.5, "w": 2.2, "h": 2.2,
         "color_key": "accent", "alpha_pct": 15},
    ],
    "section": [
        # Circle right side
        {"shape": "circle", "left": 11.0, "top": 2.2, "w": 3.0, "h": 3.0,
         "color_key": "primary", "alpha_pct": 16},
        # Ring left side
        {"shape": "ring", "left": -0.6, "top": 2.8, "w": 2.2, "h": 2.2,
         "color_key": "accent", "alpha_pct": 12},
    ],
}


def _apply_shape_alpha(shape, alpha_pct):
    """Set fill transparency on a solid-filled shape via XML alpha element."""
    sp_pr = shape._element.spPr
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is not None:
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is not None:
            # Remove existing alpha if any
            for old in srgb.findall(qn("a:alpha")):
                srgb.remove(old)
            alpha_el = srgb.makeelement(
                qn("a:alpha"), {"val": str(alpha_pct * 1000)}
            )
            srgb.append(alpha_el)


def _apply_line_alpha(shape, alpha_pct):
    """Set stroke transparency on a shape's line via XML alpha element."""
    sp_pr = shape._element.spPr
    ln = sp_pr.find(qn("a:ln"))
    if ln is not None:
        solid_fill = ln.find(qn("a:solidFill"))
        if solid_fill is not None:
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is not None:
                for old in srgb.findall(qn("a:alpha")):
                    srgb.remove(old)
                alpha_el = srgb.makeelement(
                    qn("a:alpha"), {"val": str(alpha_pct * 1000)}
                )
                srgb.append(alpha_el)


def _is_dark_background(theme):
    """Check if the theme uses a dark background (luminance < 40%)."""
    bg = theme.colors.background.lstrip("#")
    r, g, b = int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance < 0.4


def _add_bg_graphics(slide, theme, variant):
    """Add subtle translucent geometric shapes behind slide content."""
    bg_hex = theme.colors.background
    dark = _is_dark_background(theme)
    # Dark themes: use the theme color directly (no blend toward dark bg).
    # Light themes: blend partway toward background so shapes stay soft.
    blend_ratio = 0.85 if dark else 0.50

    for desc in _BG_GRAPHICS.get(variant, []):
        color_hex = getattr(theme.colors, desc["color_key"], theme.colors.primary)
        blended = _blend_color(color_hex, bg_hex, blend_ratio)
        blended_clean = blended.lstrip("#")

        left = Inches(desc["left"])
        top = Inches(desc["top"])
        width = Inches(desc["w"])
        height = Inches(desc["h"])
        shape_type = desc["shape"]

        if shape_type == "circle":
            shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(blended_clean)
            shp.line.fill.background()
            _apply_shape_alpha(shp, desc["alpha_pct"])

        elif shape_type == "ring":
            shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            shp.fill.background()  # No fill — ring only
            shp.line.color.rgb = RGBColor.from_string(blended_clean)
            shp.line.width = Pt(3.5)
            _apply_line_alpha(shp, desc["alpha_pct"])


# === Card Layouts ===

def _add_card(slide, *, left, top, width, height, theme, color_key="primary"):
    """Add a rectangle card with subtle border and left accent strip."""
    # Card background — 10% tint of theme color
    base_color = getattr(theme.colors, color_key)
    card_fill = _blend_color(base_color, theme.colors.background, 0.10)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    card.adjustments[0] = 0.05  # Corner radius
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string(card_fill.lstrip("#"))
    card.line.fill.background()  # No visible border

    # Soft drop shadow via XML
    sp_pr = card._element.spPr
    effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
    outer_shdw = effect_lst.makeelement(qn("a:outerShdw"), {
        "blurRad": "190500",   # ~15pt blur
        "dist": "63500",       # ~5pt distance
        "dir": "5400000",      # 90 degrees (straight down)
        "algn": "bl",
    })
    srgb = outer_shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgb.makeelement(qn("a:alpha"), {"val": "15000"})  # 15% opacity
    srgb.append(alpha)
    outer_shdw.append(srgb)
    effect_lst.append(outer_shdw)
    sp_pr.append(effect_lst)

    return card


# === Visual Primitives ===

# Design tokens for new visual primitives
_DESIGN_TOKENS.update({
    "badge_font_size": Pt(10),
    "badge_height": Inches(0.28),
    "progress_bar_height": Inches(0.12),
    "grid_gap": Inches(0.2),
    "numbered_badge_size": Inches(0.45),
})


def _add_tag_badge(slide, text, left, top, *, fill_color, text_color,
                   font_name, font_size=None):
    """Add a small rounded-rect pill badge with centered text."""
    if font_size is None:
        font_size = _DESIGN_TOKENS["badge_font_size"]
    badge_h = _DESIGN_TOKENS["badge_height"]
    # Width proportional to text length
    badge_w = max(Inches(0.8), Inches(len(str(text)) * 0.09 + 0.3))

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, badge_w, badge_h,
    )
    badge.adjustments[0] = 0.25  # Generous corner radius for pill shape
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(text).upper()
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = True
    if isinstance(text_color, str):
        p.font.color.rgb = RGBColor.from_string(text_color.lstrip("#"))
    else:
        p.font.color.rgb = text_color

    return badge


def _add_numbered_badge(slide, number, left, top, *, size=None,
                        fill_color, text_color, font_name):
    """Add a rounded-square badge with centered number."""
    if size is None:
        size = _DESIGN_TOKENS["numbered_badge_size"]
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size,
    )
    badge.adjustments[0] = 0.15
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(number)
    p.font.name = font_name
    p.font.size = Pt(18)
    p.font.bold = True
    if isinstance(text_color, str):
        p.font.color.rgb = RGBColor.from_string(text_color.lstrip("#"))
    else:
        p.font.color.rgb = text_color

    return badge


def _add_progress_bar(slide, left, top, width, *, percentage, track_color,
                      fill_color, label_color, font_name):
    """Add a progress bar with background track, fill bar, and percentage label."""
    bar_h = _DESIGN_TOKENS["progress_bar_height"]

    # Track (background)
    track = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, bar_h,
    )
    track.adjustments[0] = 0.5
    track.fill.solid()
    track.fill.fore_color.rgb = RGBColor.from_string(track_color.lstrip("#"))
    track.line.fill.background()

    # Fill (proportional)
    pct = max(0.0, min(1.0, percentage / 100.0))
    fill_w = max(Inches(0.1), int(width * pct))
    fill_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, fill_w, bar_h,
    )
    fill_bar.adjustments[0] = 0.5
    fill_bar.fill.solid()
    fill_bar.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
    fill_bar.line.fill.background()

    # Label
    label_w = Inches(0.6)
    _add_text_box(slide, f"{int(percentage)}%",
                  left=left + width + Inches(0.1), top=top - Inches(0.02),
                  width=label_w, height=Inches(0.2),
                  font_name=font_name, font_size=_DESIGN_TOKENS["badge_font_size"],
                  font_color=label_color, bold=True, parse_markdown=False)


def _add_thin_divider(slide, left, top, height, *, color):
    """Add a thin vertical line divider."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.02), height,
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    line.line.fill.background()
    return line


def _add_horizontal_rule(slide, left, top, width, *, color):
    """Add a thin horizontal accent underline."""
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    rule.line.fill.background()
    return rule


# === Chart Rendering ===

def _build_chart_palette(theme):
    """Build a 6-color chart palette from theme colors."""
    colors = [theme.colors.accent, theme.colors.primary, theme.colors.secondary]
    # Add 40% tints toward white
    for c in [theme.colors.accent, theme.colors.primary, theme.colors.secondary]:
        colors.append(_blend_color(c, "#FFFFFF", 0.60))
    return colors


def _style_chart(chart, spec, theme):
    """Apply theme styling to a python-pptx chart object."""
    try:
        chart.chart_style = 2  # Minimal style
        plot = chart.plots[0]
        plot.format.fill.background()  # Transparent plot area
    except Exception:
        pass

    # Series colors
    palette = _build_chart_palette(theme)
    is_line_chart = spec.chart_type.value in ("line",)
    try:
        for idx, series in enumerate(chart.series):
            color_hex = palette[idx % len(palette)].lstrip("#")
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(color_hex)
            # Thicker lines for line charts
            if is_line_chart:
                series.format.line.width = Pt(3)
                series.format.line.color.rgb = RGBColor.from_string(color_hex)
    except Exception:
        pass

    # Data labels on first series for line charts
    if is_line_chart:
        try:
            chart.series[0].has_data_labels = True
            dl = chart.series[0].data_labels
            dl.font.size = _DESIGN_TOKENS["code_size"]
            dl.font.color.rgb = RGBColor.from_string(
                theme.colors.text_light.lstrip("#")
            )
        except Exception:
            pass

    # Axis labels
    try:
        if hasattr(chart, 'value_axis') and chart.value_axis is not None:
            va = chart.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = RGBColor.from_string(
                theme.colors.text_light.lstrip("#")
            )
            va.format.line.fill.background()  # Hide redundant axis line
            if va.has_tick_labels:
                va.tick_labels.font.size = _DESIGN_TOKENS["body_small_size"]
                va.tick_labels.font.color.rgb = RGBColor.from_string(
                    theme.colors.text_light.lstrip("#")
                )
    except Exception:
        pass

    try:
        if hasattr(chart, 'category_axis') and chart.category_axis is not None:
            ca = chart.category_axis
            ca.has_major_gridlines = False
            if ca.has_tick_labels:
                ca.tick_labels.font.size = _DESIGN_TOKENS["body_small_size"]
                ca.tick_labels.font.color.rgb = RGBColor.from_string(
                    theme.colors.text_light.lstrip("#")
                )
    except Exception:
        pass

    # Legend
    try:
        if chart.has_legend and chart.legend is not None:
            chart.legend.font.size = _DESIGN_TOKENS["code_size"]
    except Exception:
        pass


def _add_chart(slide, spec, theme):
    """Add a native python-pptx chart to the slide. Returns True on success."""
    if not _CHARTS_AVAILABLE:
        return False
    try:
        from core.schemas.studio_schema import ChartType

        # Chart container card — rounded rect behind chart area
        card_margin = Inches(0.15)
        _add_card(slide,
                  left=CHART_LEFT - card_margin,
                  top=CHART_TOP - card_margin,
                  width=CHART_WIDTH + card_margin * 2,
                  height=CHART_HEIGHT + card_margin * 2,
                  theme=theme, color_key="primary")

        if spec.chart_type.value == "scatter":
            chart_data = XyChartData()
            series = chart_data.add_series(spec.title or "Data")
            for pt in spec.points:
                series.add_data_point(pt.x, pt.y)
            chart_type = XL_CHART_TYPE.XY_SCATTER
        else:
            chart_data = CategoryChartData()
            chart_data.categories = spec.categories
            series_list = spec.series
            if spec.chart_type == ChartType.pie and len(series_list) > 1:
                series_list = series_list[:1]
            for s in series_list:
                chart_data.add_series(s.name, s.values)
            chart_type = _CHART_TYPE_MAP.get(spec.chart_type.value, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_frame = slide.shapes.add_chart(
            chart_type, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = len(spec.series) > 1 or spec.chart_type.value == "pie"

        _style_chart(chart, spec, theme)

        # Enable data labels on bar/pie charts with <=6 categories
        cat_count = len(spec.categories) if spec.categories else 0
        if spec.chart_type.value in ("bar", "pie") and cat_count <= 6:
            try:
                for series_obj in chart.series:
                    series_obj.has_data_labels = True
                    dl = series_obj.data_labels
                    dl.font.size = _DESIGN_TOKENS["code_size"]
                    dl.font.color.rgb = RGBColor.from_string(
                        theme.colors.text.lstrip("#")
                    )
            except Exception:
                pass

        return True
    except Exception:
        return False


# === Slide-Type Renderer Functions ===

def _render_title(slide, slide_data, theme, **kwargs):
    """Title slide: centered title + subtitle, with optional dark background."""
    title_bg = getattr(theme.colors, "title_background", None)
    if title_bg:
        # Dark title slide — override background and use white text
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(title_bg.lstrip("#"))
        title_color = "#FFFFFF"
        subtitle_color = "#CCCCCC"
    else:
        _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))
        title_color = theme.colors.primary
        subtitle_color = theme.colors.text_light

    # Background geometric graphics (before text for z-order)
    slide_index = kwargs.get("slide_index", 0)
    total_slides = kwargs.get("total_slides", 1)
    is_closing = slide_index == total_slides - 1 and total_slides > 1
    _add_bg_graphics(slide, theme, variant="closing" if is_closing else "title")

    # Check for metadata-driven enhancements
    metadata = getattr(slide_data, "metadata", None) or {}
    if isinstance(metadata, dict):
        badge_text = metadata.get("badge")
        meta_date = metadata.get("date")
        meta_category = metadata.get("category")
    else:
        badge_text = None
        meta_date = None
        meta_category = None

    # Multi-color title: if title contains \n, render line 1 in primary, line 2 in accent
    title_text = slide_data.title or ""
    title_len = len(title_text)

    # Adapt font size and spacing for long titles
    if title_len > 60:
        title_font_size = Pt(36)
        title_top = Inches(2.2)
        title_height = Inches(2.0)
        subtitle_top = Inches(4.4)
    elif title_len > 40:
        title_font_size = Pt(40)
        title_top = Inches(2.3)
        title_height = Inches(1.8)
        subtitle_top = Inches(4.3)
    else:
        title_font_size = _DESIGN_TOKENS["title_size"]
        title_top = Inches(2.5)
        title_height = Inches(1.5)
        subtitle_top = Inches(4.2)

    if "\n" in title_text:
        lines = title_text.split("\n", 1)
        _add_text_box(slide, lines[0],
                      left=MARGIN_LEFT, top=title_top,
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_heading, font_size=title_font_size,
                      font_color=title_color, alignment=PP_ALIGN.CENTER,
                      bold=True)
        _add_text_box(slide, lines[1],
                      left=MARGIN_LEFT, top=title_top + Inches(0.8),
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_heading, font_size=title_font_size,
                      font_color=theme.colors.accent, alignment=PP_ALIGN.CENTER,
                      bold=True)
    else:
        _add_text_box(slide, title_text,
                      left=MARGIN_LEFT, top=title_top,
                      width=CONTENT_WIDTH, height=title_height,
                      font_name=theme.font_heading, font_size=title_font_size,
                      font_color=title_color, alignment=PP_ALIGN.CENTER,
                      bold=True)

    subtitle_el = _find_element(slide_data, "subtitle")
    if subtitle_el and subtitle_el.content:
        _add_text_box(slide, subtitle_el.content,
                      left=MARGIN_LEFT, top=subtitle_top,
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["subheading_size"],
                      font_color=subtitle_color, alignment=PP_ALIGN.CENTER)

    # Closing title enhancements: stats with dividers
    slide_index = kwargs.get("slide_index", 0)
    total_slides = kwargs.get("total_slides", 1)
    stat_el = _find_element(slide_data, "stat_callout")
    if slide_index == total_slides - 1 and stat_el:
        # Parse stat data
        stats = []
        if isinstance(stat_el.content, list):
            for item in stat_el.content[:3]:
                if isinstance(item, dict) and "value" in item:
                    stats.append(item)
        if stats:
            col_count = len(stats)
            col_width = CONTENT_WIDTH / col_count
            stat_top = Inches(5.2)
            for idx, stat in enumerate(stats):
                col_left = MARGIN_LEFT + col_width * idx
                _add_text_box(slide, stat.get("value", ""),
                              left=col_left, top=stat_top,
                              width=col_width, height=Inches(0.6),
                              font_name=theme.font_heading,
                              font_size=Pt(28),
                              font_color=theme.colors.accent,
                              alignment=PP_ALIGN.CENTER, bold=True)
                _add_text_box(slide, stat.get("label", ""),
                              left=col_left, top=stat_top + Inches(0.55),
                              width=col_width, height=Inches(0.4),
                              font_name=theme.font_body,
                              font_size=_DESIGN_TOKENS["footer_size"],
                              font_color=subtitle_color,
                              alignment=PP_ALIGN.CENTER)
                # Thin divider between stat columns
                if idx < col_count - 1:
                    _add_thin_divider(slide,
                                      left=col_left + col_width - Inches(0.01),
                                      top=stat_top,
                                      height=Inches(0.9),
                                      color=theme.colors.text_light)

    # Metadata footer (date, category)
    if meta_date or meta_category:
        footer_parts = []
        if meta_date:
            footer_parts.append(str(meta_date))
        if meta_category:
            footer_parts.append(str(meta_category))
        footer_text = "  |  ".join(footer_parts)
        _add_text_box(slide, footer_text,
                      left=MARGIN_LEFT, top=Inches(6.5),
                      width=CONTENT_WIDTH, height=Inches(0.4),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["footer_size"],
                      font_color=subtitle_color,
                      alignment=PP_ALIGN.CENTER)


def _render_content(slide, slide_data, theme, **kwargs):
    """Standard content slide: title + body/bullets with varied card styles."""
    slide_index = kwargs.get("slide_index", 0)
    card_style = slide_index % 3

    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_el = _find_element(slide_data, "body")
    bullet_el = _find_element(slide_data, "bullet_list")

    content_left = BODY_LEFT
    content_width = BODY_WIDTH

    if card_style == 0:
        # Style 0: Full card with left accent strip (original)
        _add_card(slide, left=BODY_LEFT, top=body_top,
                  width=BODY_WIDTH, height=body_height,
                  theme=theme, color_key="primary")
        content_left = BODY_LEFT + Inches(0.12)
        content_width = BODY_WIDTH - Inches(0.12)
    elif card_style == 1:
        # Style 1: Top accent bar above content area
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, body_top, BODY_WIDTH, Inches(0.08),
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.accent.lstrip("#")
        )
        top_bar.line.fill.background()
        # Offset content below the bar (0.08" bar + 0.15" gap)
        body_top = body_top + Inches(0.23)
        body_height = body_height - Inches(0.23)
    else:
        # Style 2: Horizontal rule below title
        rule_top = body_top - Inches(0.15)
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            BODY_LEFT, rule_top, BODY_WIDTH, Inches(0.03),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.accent.lstrip("#")
        )
        rule.line.fill.background()

    if bullet_el and isinstance(bullet_el.content, list):
        _add_icon_bullet_list(slide, bullet_el.content,
                              left=content_left, top=body_top,
                              width=content_width, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=content_left, top=body_top,
                      width=content_width, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_two_column(slide, slide_data, theme, **kwargs):
    """Two-column layout: title + left/right body areas."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_elements = [el for el in slide_data.elements if el.type == "body"]
    bullet_elements = [el for el in slide_data.elements if el.type == "bullet_list"]

    # Left column
    left_content = body_elements[0].content if body_elements else ""
    if bullet_elements and isinstance(bullet_elements[0].content, list):
        _add_icon_bullet_list(slide, bullet_elements[0].content,
                              left=MARGIN_LEFT, top=body_top,
                              width=COLUMN_WIDTH, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    else:
        _add_text_box(slide, left_content,
                      left=MARGIN_LEFT, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    # Right column
    right_content = body_elements[1].content if len(body_elements) > 1 else ""
    if len(bullet_elements) > 1 and isinstance(bullet_elements[1].content, list):
        _add_icon_bullet_list(slide, bullet_elements[1].content,
                              left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                              width=COLUMN_WIDTH, height=body_height,
                              font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                              font_color=theme.colors.text,
                              icon_color=theme.colors.accent,
                              accent_color=theme.colors.accent)
    else:
        _add_text_box(slide, right_content,
                      left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_comparison(slide, slide_data, theme, **kwargs):
    """Comparison slide: title + two labeled columns in cards with optional callout."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))

    # Check for callout_box element
    callout_el = _find_element(slide_data, "callout_box")
    has_callout = callout_el and callout_el.content
    callout_height = Inches(0.9) if has_callout else Inches(0)
    callout_gap = Inches(0.15) if has_callout else Inches(0)

    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM - callout_height - callout_gap
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    body_elements = [el for el in slide_data.elements if el.type == "body"]
    left_text = body_elements[0].content if body_elements else ""
    right_text = body_elements[1].content if len(body_elements) > 1 else ""

    # Column header labels
    left_label = ""
    right_label = ""
    if body_elements:
        # Extract label from first line if it's short
        if isinstance(left_text, str) and "\n" in left_text:
            first_line, rest = left_text.split("\n", 1)
            if len(first_line) < 30:
                left_label = first_line
                left_text = rest
    if len(body_elements) > 1:
        if isinstance(right_text, str) and "\n" in right_text:
            first_line, rest = right_text.split("\n", 1)
            if len(first_line) < 30:
                right_label = first_line
                right_text = rest

    # Labeled column headers above cards
    header_height = Inches(0.35) if (left_label or right_label) else Inches(0)
    header_gap = Inches(0.05) if header_height else Inches(0)
    card_top = body_top + header_height + header_gap
    card_height = body_height - header_height - header_gap

    if left_label:
        _add_text_box(slide, left_label,
                      left=MARGIN_LEFT, top=body_top,
                      width=COLUMN_WIDTH, height=header_height,
                      font_name=theme.font_heading,
                      font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.secondary, bold=True)
        _add_horizontal_rule(slide, left=MARGIN_LEFT, top=body_top + header_height,
                             width=COLUMN_WIDTH, color=theme.colors.secondary)

    if right_label:
        right_left = MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP
        _add_text_box(slide, right_label,
                      left=right_left, top=body_top,
                      width=COLUMN_WIDTH, height=header_height,
                      font_name=theme.font_heading,
                      font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.accent, bold=True)
        _add_horizontal_rule(slide, left=right_left, top=body_top + header_height,
                             width=COLUMN_WIDTH, color=theme.colors.accent)

    # Left card (primary)
    _add_card(slide, left=MARGIN_LEFT, top=card_top,
              width=COLUMN_WIDTH, height=card_height,
              theme=theme, color_key="primary")
    _add_text_box(slide, left_text,
                  left=MARGIN_LEFT + Inches(0.12), top=card_top,
                  width=COLUMN_WIDTH - Inches(0.12), height=card_height,
                  font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                  font_color=theme.colors.text,
                  accent_color=theme.colors.accent)

    # Right card (secondary)
    right_left = MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP
    _add_card(slide, left=right_left, top=card_top,
              width=COLUMN_WIDTH, height=card_height,
              theme=theme, color_key="secondary")
    _add_text_box(slide, right_text,
                  left=right_left + Inches(0.12), top=card_top,
                  width=COLUMN_WIDTH - Inches(0.12), height=card_height,
                  font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                  font_color=theme.colors.text,
                  accent_color=theme.colors.accent)

    # Callout box at bottom
    if has_callout:
        callout_data = callout_el.content
        if isinstance(callout_data, str):
            try:
                import json
                callout_data = json.loads(callout_data)
            except (json.JSONDecodeError, TypeError):
                callout_data = {"text": callout_data}
        if isinstance(callout_data, dict):
            callout_text = callout_data.get("text", "")
            attribution = callout_data.get("attribution", "")
        else:
            callout_text = str(callout_data)
            attribution = ""

        callout_top = card_top + card_height + callout_gap
        # Dark rounded rect
        callout_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN_LEFT, callout_top, CONTENT_WIDTH, callout_height,
        )
        callout_bg.adjustments[0] = 0.08
        callout_fill = _blend_color(theme.colors.primary, theme.colors.background, 0.85)
        callout_bg.fill.solid()
        callout_bg.fill.fore_color.rgb = RGBColor.from_string(callout_fill.lstrip("#"))
        callout_bg.line.fill.background()

        display_text = f"\u201C{callout_text}\u201D"
        if attribution:
            display_text += f"  \u2014 {attribution}"
        _add_text_box(slide, display_text,
                      left=MARGIN_LEFT + Inches(0.3), top=callout_top + Inches(0.1),
                      width=CONTENT_WIDTH - Inches(0.6),
                      height=callout_height - Inches(0.2),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text_light)

    _render_takeaway(slide, slide_data, theme)
    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_timeline(slide, slide_data, theme, **kwargs):
    """Timeline/roadmap slide: card-based or circle+bullet layout."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    if _has_takeaway(slide_data):
        body_height -= _TAKEAWAY_RESERVE

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    bullet_el = _find_element(slide_data, "bullet_list")
    body_el = _find_element(slide_data, "body")

    if bullet_el and isinstance(bullet_el.content, list):
        items = bullet_el.content
        # Detect pipe-delimited format: "Date | Title | Description | TAG"
        has_pipe_format = any("|" in str(item) for item in items)

        if has_pipe_format:
            # Card-based timeline layout
            card_gap = Inches(0.15)
            card_height = min(
                (body_height - card_gap * (len(items) - 1)) / len(items),
                Inches(1.3),
            )

            for idx, item in enumerate(items):
                parts = [p.strip() for p in str(item).split("|")]
                date_text = parts[0] if len(parts) > 0 else ""
                title_text = parts[1] if len(parts) > 1 else ""
                desc_text = parts[2] if len(parts) > 2 else ""
                tag_text = parts[3] if len(parts) > 3 else ""

                card_top = body_top + idx * (card_height + card_gap)

                # Card background
                _add_card(slide, left=BODY_LEFT, top=card_top,
                          width=BODY_WIDTH, height=card_height,
                          theme=theme, color_key="primary")

                # Date label in accent color
                _add_text_box(slide, date_text,
                              left=BODY_LEFT + Inches(0.2), top=card_top + Inches(0.08),
                              width=Inches(1.5), height=Inches(0.3),
                              font_name=theme.font_body,
                              font_size=_DESIGN_TOKENS["body_small_size"],
                              font_color=theme.colors.accent,
                              bold=True, parse_markdown=False)

                # Bold title
                _add_text_box(slide, title_text,
                              left=BODY_LEFT + Inches(1.8), top=card_top + Inches(0.08),
                              width=BODY_WIDTH - Inches(3.5), height=Inches(0.3),
                              font_name=theme.font_heading,
                              font_size=_DESIGN_TOKENS["body_size"],
                              font_color=theme.colors.text,
                              bold=True)

                # Description
                if desc_text:
                    _add_text_box(slide, desc_text,
                                  left=BODY_LEFT + Inches(1.8),
                                  top=card_top + Inches(0.4),
                                  width=BODY_WIDTH - Inches(3.5),
                                  height=card_height - Inches(0.5),
                                  font_name=theme.font_body,
                                  font_size=_DESIGN_TOKENS["body_small_size"],
                                  font_color=theme.colors.text_light)

                # Tag badge
                if tag_text:
                    _add_tag_badge(slide, tag_text,
                                   left=BODY_LEFT + BODY_WIDTH - Inches(1.5),
                                   top=card_top + Inches(0.1),
                                   fill_color=theme.colors.secondary,
                                   text_color="#FFFFFF",
                                   font_name=theme.font_body)
        else:
            # Classic circle+bullet layout (backward compatible)
            item_count = len(items)
            for idx in range(item_count):
                circle_top = body_top + Inches(idx * 0.55 + 0.1)
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    BODY_LEFT, circle_top, Inches(0.25), Inches(0.25),
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor.from_string(
                    theme.colors.accent.lstrip("#")
                )
                circle.line.fill.background()

            _add_bullet_list(slide, items,
                             left=BODY_LEFT + Inches(0.4), top=body_top,
                             width=BODY_WIDTH - Inches(0.4), height=body_height,
                             font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                             font_color=theme.colors.text,
                             accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        # Handle list-typed body as pipe-delimited timeline items
        if isinstance(body_el.content, list):
            items = body_el.content
            has_pipe_format = any("|" in str(item) for item in items)
            if has_pipe_format:
                card_gap = Inches(0.15)
                card_height = min(
                    (body_height - card_gap * (len(items) - 1)) / len(items),
                    Inches(1.3),
                )
                for idx, item in enumerate(items):
                    parts = [p.strip() for p in str(item).split("|")]
                    date_text = parts[0] if len(parts) > 0 else ""
                    title_text = parts[1] if len(parts) > 1 else ""
                    desc_text = parts[2] if len(parts) > 2 else ""

                    card_top = body_top + idx * (card_height + card_gap)
                    _add_card(slide, left=BODY_LEFT, top=card_top,
                              width=BODY_WIDTH, height=card_height,
                              theme=theme, color_key="primary")
                    _add_text_box(slide, date_text,
                                  left=BODY_LEFT + Inches(0.2), top=card_top + Inches(0.08),
                                  width=Inches(1.5), height=Inches(0.3),
                                  font_name=theme.font_body,
                                  font_size=_DESIGN_TOKENS["body_small_size"],
                                  font_color=theme.colors.accent,
                                  bold=True, parse_markdown=False)
                    _add_text_box(slide, title_text,
                                  left=BODY_LEFT + Inches(1.8), top=card_top + Inches(0.08),
                                  width=BODY_WIDTH - Inches(3.5), height=Inches(0.3),
                                  font_name=theme.font_heading,
                                  font_size=_DESIGN_TOKENS["body_size"],
                                  font_color=theme.colors.text,
                                  bold=True)
                    if desc_text:
                        _add_text_box(slide, desc_text,
                                      left=BODY_LEFT + Inches(1.8),
                                      top=card_top + Inches(0.4),
                                      width=BODY_WIDTH - Inches(3.5),
                                      height=card_height - Inches(0.5),
                                      font_name=theme.font_body,
                                      font_size=_DESIGN_TOKENS["body_small_size"],
                                      font_color=theme.colors.text_light)
            else:
                _add_bullet_list(slide, items,
                                 left=BODY_LEFT, top=body_top,
                                 width=BODY_WIDTH, height=body_height,
                                 font_name=theme.font_body,
                                 font_size=_DESIGN_TOKENS["body_size"],
                                 font_color=theme.colors.text,
                                 accent_color=theme.colors.accent)
        else:
            _add_text_box(slide, body_el.content,
                          left=BODY_LEFT, top=body_top,
                          width=BODY_WIDTH, height=body_height,
                          font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                          font_color=theme.colors.text,
                          accent_color=theme.colors.accent)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_chart(slide, slide_data, theme, **kwargs):
    """Chart slide: title + native chart or text fallback."""
    from core.studio.slides.charts import parse_chart_spec, normalize_chart_spec

    _render_kicker(slide, slide_data, theme)

    _, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    chart_el = _find_element(slide_data, "chart")
    body_el = _find_element(slide_data, "body")
    chart_rendered = False

    if chart_el and chart_el.content:
        spec = parse_chart_spec(chart_el.content)
        if spec is not None:
            spec = normalize_chart_spec(spec)
            chart_rendered = _add_chart(slide, spec, theme)

    # Fallback: text placeholder if chart parse/render failed
    if not chart_rendered and chart_el and chart_el.content:
        content = chart_el.content if isinstance(chart_el.content, str) else str(chart_el.content)
        _add_text_box(slide, f"[Chart: {content}]",
                      left=BODY_LEFT, top=BODY_TOP,
                      width=BODY_WIDTH, height=Inches(3.0),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.secondary,
                      alignment=PP_ALIGN.CENTER)

    # Body/caption below chart
    if body_el and body_el.content:
        caption_top = CAPTION_TOP if chart_rendered else Inches(5.0)
        caption_h = CAPTION_HEIGHT
        if _has_takeaway(slide_data):
            caption_h = Inches(0.3)  # Shrink to leave room for takeaway
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=caption_top,
                      width=BODY_WIDTH, height=caption_h,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_image_text(slide, slide_data, theme, **kwargs):
    """Image+text slide: split layout with image (or placeholder) and body."""
    body_top, title_font, title_top = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    image_el = _find_element(slide_data, "image")
    body_el = _find_element(slide_data, "body")
    has_body = body_el and body_el.content

    # Check for a generated image or external URL
    images = kwargs.get("images") or {}
    img_buf = images.get(slide_data.id) if slide_data.id else None

    # Try external URL if no generated image
    if img_buf is None and image_el and image_el.content:
        ext_url = extract_image_url(image_el.content)
        if ext_url:
            img_buf = _download_image_url(ext_url)

    # Use full width when there's no body text
    img_width = COLUMN_WIDTH if has_body else BODY_WIDTH

    if img_buf is not None:
        # Embed actual image
        img_buf.seek(0)
        pic = slide.shapes.add_picture(
            img_buf, MARGIN_LEFT, body_top, img_width, body_height,
        )
        # Add 2pt border in primary color
        pic.line.width = Pt(2)
        pic.line.color.rgb = RGBColor.from_string(
            theme.colors.primary.lstrip("#")
        )
    else:
        # Text placeholder fallback
        placeholder_text = "[Image]"
        if image_el and image_el.content:
            placeholder_text = f"[Image: {image_el.content}]"
        _add_text_box(slide, placeholder_text,
                      left=MARGIN_LEFT, top=body_top,
                      width=img_width, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.secondary,
                      alignment=PP_ALIGN.CENTER)

    if has_body:
        _add_text_box(slide, body_el.content,
                      left=MARGIN_LEFT + COLUMN_WIDTH + COLUMN_GAP, top=body_top,
                      width=COLUMN_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_quote(slide, slide_data, theme, **kwargs):
    """Quote slide: large quote text with accent bar and attribution."""
    quote_el = _find_element(slide_data, "quote")
    body_el = _find_element(slide_data, "body")

    quote_text = ""
    if quote_el and quote_el.content:
        raw = quote_el.content.strip().strip('"\u201C\u201D')
        quote_text = f"\u201C{raw}\u201D"
    elif slide_data.title:
        quote_text = slide_data.title

    # Vertical accent bar to the left of quote
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.2), Inches(1.5), Inches(0.12), Inches(4.0),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor.from_string(
        theme.colors.accent.lstrip("#")
    )
    accent_bar.line.fill.background()

    _add_text_box(slide, quote_text,
                  left=Inches(1.5), top=Inches(1.5),
                  width=Inches(10.333), height=Inches(4.0),
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["quote_size"],
                  font_color=theme.colors.primary,
                  alignment=PP_ALIGN.CENTER)

    if body_el and body_el.content:
        _add_text_box(slide, f"\u2014 {body_el.content}",
                      left=Inches(1.5), top=Inches(5.8),
                      width=Inches(10.333), height=Inches(0.8),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["attribution_size"],
                      font_color=theme.colors.text_light,
                      alignment=PP_ALIGN.CENTER)

    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_code(slide, slide_data, theme, **kwargs):
    """Code slide: title + monospace code block."""
    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=TITLE_TOP,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["heading_size"],
                  font_color=theme.colors.primary, bold=True)

    code_el = _find_element(slide_data, "code")
    if code_el and code_el.content:
        _add_text_box(slide, code_el.content,
                      left=Inches(1.0), top=BODY_TOP,
                      width=Inches(11.333), height=BODY_HEIGHT,
                      font_name="Courier New", font_size=_DESIGN_TOKENS["code_size"],
                      font_color=theme.colors.text, parse_markdown=False)

    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_team(slide, slide_data, theme, **kwargs):
    """Team/credits slide: title + team member list."""
    body_top, title_font, title_top = _compute_body_top(slide_data.title)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    bullet_el = _find_element(slide_data, "bullet_list")
    body_el = _find_element(slide_data, "body")

    if bullet_el and isinstance(bullet_el.content, list):
        _add_bullet_list(slide, bullet_el.content,
                         left=BODY_LEFT, top=body_top,
                         width=BODY_WIDTH, height=body_height,
                         font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                         font_color=theme.colors.text,
                         accent_color=theme.colors.accent)
    elif body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=body_top,
                      width=BODY_WIDTH, height=body_height,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text,
                      accent_color=theme.colors.accent)

    _add_decorative_accent(slide, theme, kwargs.get("slide_index", 0))
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_stat(slide, slide_data, theme, **kwargs):
    """Stat slide: 1-3 large stat callouts with labels."""
    _render_kicker(slide, slide_data, theme)

    _, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    stat_el = _find_element(slide_data, "stat_callout")
    body_el = _find_element(slide_data, "body")

    stats = []
    if stat_el and isinstance(stat_el.content, list):
        for item in stat_el.content[:3]:
            if isinstance(item, dict) and "value" in item:
                stats.append(item)
    elif stat_el and isinstance(stat_el.content, str):
        # Fallback: try parsing JSON string
        try:
            import json
            parsed = json.loads(stat_el.content)
            if isinstance(parsed, list):
                for item in parsed[:3]:
                    if isinstance(item, dict) and "value" in item:
                        stats.append(item)
        except (json.JSONDecodeError, TypeError):
            pass

    if stats:
        col_count = len(stats)
        col_width = CONTENT_WIDTH / col_count
        for idx, stat in enumerate(stats):
            col_left = MARGIN_LEFT + col_width * idx
            # Large value
            _add_text_box(slide, stat.get("value", ""),
                          left=col_left, top=Inches(2.5),
                          width=col_width, height=Inches(2.0),
                          font_name=theme.font_heading,
                          font_size=_DESIGN_TOKENS["stat_callout_size"],
                          font_color=theme.colors.accent,
                          alignment=PP_ALIGN.CENTER, bold=True)
            # Label below
            _add_text_box(slide, stat.get("label", ""),
                          left=col_left, top=Inches(4.5),
                          width=col_width, height=Inches(0.8),
                          font_name=theme.font_body,
                          font_size=_DESIGN_TOKENS["stat_label_size"],
                          font_color=theme.colors.text_light,
                          alignment=PP_ALIGN.CENTER)
    elif body_el and body_el.content:
        # Graceful fallback: render as body text
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=BODY_TOP,
                      width=BODY_WIDTH, height=BODY_HEIGHT,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text)

    # Optional body/context below stats
    if stats and body_el and body_el.content:
        stat_body_h = Inches(1.2)
        if _has_takeaway(slide_data):
            stat_body_h = Inches(0.5)  # Shrink to leave room for takeaway
        _add_text_box(slide, body_el.content,
                      left=BODY_LEFT, top=Inches(5.5),
                      width=BODY_WIDTH, height=stat_body_h,
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.text_light,
                      alignment=PP_ALIGN.CENTER)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _render_image_full(slide, slide_data, theme, **kwargs):
    """Full-bleed image slide: image covers entire slide with dark overlay and white text."""
    images = kwargs.get("images") or {}
    img_buf = images.get(slide_data.id) if slide_data.id else None

    # Try external URL if no generated image
    if img_buf is None:
        image_el = _find_element(slide_data, "image")
        if image_el and image_el.content:
            ext_url = extract_image_url(image_el.content)
            if ext_url:
                img_buf = _download_image_url(ext_url)

    if img_buf is not None:
        # Embed full-bleed image covering entire slide
        img_buf.seek(0)
        slide.shapes.add_picture(
            img_buf, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
        )
    else:
        # Dark background fallback when no image available
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string("1A1A2E")

    # Gradient dark overlay: transparent at top, dark at bottom (where text sits)
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    overlay.fill.gradient()
    overlay.line.fill.background()

    # Configure gradient stops via XML for precise alpha control
    sp_pr = overlay._element.spPr
    grad_fill = sp_pr.find(qn("a:gradFill"))
    if grad_fill is not None:
        # Set linear angle: 16200000 = 270 degrees (bottom-to-top in EMU)
        lin = grad_fill.find(qn("a:lin"))
        if lin is None:
            lin = grad_fill.makeelement(qn("a:lin"), {})
            grad_fill.append(lin)
        lin.set("ang", "16200000")
        lin.set("scaled", "1")

        # Replace auto-generated stops with our custom ones
        gs_lst = grad_fill.find(qn("a:gsLst"))
        if gs_lst is not None:
            for gs in list(gs_lst):
                gs_lst.remove(gs)

            # Stop 0 (top of slide): black at 0% opacity (fully transparent)
            gs0 = gs_lst.makeelement(qn("a:gs"), {"pos": "0"})
            srgb0 = gs0.makeelement(qn("a:srgbClr"), {"val": "000000"})
            alpha0 = srgb0.makeelement(qn("a:alpha"), {"val": "0"})
            srgb0.append(alpha0)
            gs0.append(srgb0)
            gs_lst.append(gs0)

            # Stop 1 (bottom of slide): black at 80% opacity
            gs1 = gs_lst.makeelement(qn("a:gs"), {"pos": "100000"})
            srgb1 = gs1.makeelement(qn("a:srgbClr"), {"val": "000000"})
            alpha1 = srgb1.makeelement(qn("a:alpha"), {"val": "80000"})
            srgb1.append(alpha1)
            gs1.append(srgb1)
            gs_lst.append(gs1)

    # Title — centered, white, large
    _add_text_box(slide, slide_data.title or "",
                  left=Inches(1.5), top=Inches(2.0),
                  width=Inches(10.333), height=Inches(2.0),
                  font_name=theme.font_heading, font_size=_DESIGN_TOKENS["title_size"],
                  font_color="#FFFFFF", alignment=PP_ALIGN.CENTER,
                  bold=True)

    # Optional body text below title
    body_el = _find_element(slide_data, "body")
    if body_el and body_el.content:
        _add_text_box(slide, body_el.content,
                      left=Inches(2.0), top=Inches(4.5),
                      width=Inches(9.333), height=Inches(2.0),
                      font_name=theme.font_body, font_size=_DESIGN_TOKENS["subheading_size"],
                      font_color="#DDDDDD", alignment=PP_ALIGN.CENTER)


def _render_section_divider(slide, slide_data, theme, **kwargs):
    """Section divider slide: large section number + title, whitespace-heavy."""
    section_number = kwargs.get("section_number", 1)

    # Use dark background like title slide
    title_bg = getattr(theme.colors, "title_background", None)
    if title_bg:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(title_bg.lstrip("#"))
        title_color = "#FFFFFF"
        subtitle_color = "#CCCCCC"
    else:
        _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))
        title_color = theme.colors.primary
        subtitle_color = theme.colors.text_light

    # Background geometric graphics (before text for z-order)
    _add_bg_graphics(slide, theme, variant="section")

    # Large section number in accent color
    section_num = str(section_number)
    _add_text_box(slide, section_num,
                  left=MARGIN_LEFT, top=Inches(1.5),
                  width=CONTENT_WIDTH, height=Inches(2.0),
                  font_name=theme.font_heading,
                  font_size=Pt(120),
                  font_color=theme.colors.accent,
                  alignment=PP_ALIGN.CENTER,
                  bold=True, parse_markdown=False)

    # Section title
    _add_text_box(slide, slide_data.title or "",
                  left=MARGIN_LEFT, top=Inches(3.8),
                  width=CONTENT_WIDTH, height=Inches(1.5),
                  font_name=theme.font_heading,
                  font_size=_DESIGN_TOKENS["heading_size"],
                  font_color=title_color,
                  alignment=PP_ALIGN.CENTER,
                  bold=True)

    # Accent underline bar below the title
    _add_horizontal_rule(slide,
                         left=MARGIN_LEFT + (CONTENT_WIDTH - Inches(3.0)) / 2,
                         top=Inches(5.15),
                         width=Inches(3.0),
                         color=theme.colors.accent)

    # Optional subtitle
    subtitle_el = _find_element(slide_data, "subtitle")
    if subtitle_el and subtitle_el.content:
        _add_text_box(slide, subtitle_el.content,
                      left=MARGIN_LEFT, top=Inches(5.3),
                      width=CONTENT_WIDTH, height=Inches(0.8),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["body_size"],
                      font_color=subtitle_color,
                      alignment=PP_ALIGN.CENTER)

    # Decorative dots in bottom-right corner
    dot_colors = [theme.colors.accent, theme.colors.primary, theme.colors.text_light]
    for i, dot_color in enumerate(dot_colors):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SLIDE_WIDTH - Inches(1.2 + i * 0.35), SLIDE_HEIGHT - Inches(0.6),
            Inches(0.18), Inches(0.18),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor.from_string(dot_color.lstrip("#"))
        dot.line.fill.background()


def _render_agenda(slide, slide_data, theme, **kwargs):
    """Agenda slide: numbered card grid of section items."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    bullet_el = _find_element(slide_data, "bullet_list")
    items = []
    if bullet_el and isinstance(bullet_el.content, list):
        items = bullet_el.content[:9]  # Max 9 items

    if not items:
        _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))
        return

    # Grid layout: 2 columns for <=6 items, 3 columns for >6
    cols = 3 if len(items) > 6 else 2
    rows = (len(items) + cols - 1) // cols
    grid_gap = _DESIGN_TOKENS["grid_gap"]

    card_width = (CONTENT_WIDTH - grid_gap * (cols - 1)) / cols
    available_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM - Inches(0.2)
    card_height = (available_height - grid_gap * (rows - 1)) / rows
    card_height = min(card_height, Inches(2.0))

    badge_size = _DESIGN_TOKENS["numbered_badge_size"]

    for idx, item in enumerate(items):
        item = str(item)
        col = idx % cols
        row = idx // cols
        card_left = MARGIN_LEFT + col * (card_width + grid_gap)
        card_top = body_top + row * (card_height + grid_gap)

        # Parse "Title: Description" format
        if ":" in item:
            parts = item.split(":", 1)
            item_title = parts[0].strip()
            item_desc = parts[1].strip()
        else:
            item_title = item
            item_desc = ""

        # Card background
        _add_card(slide, left=card_left, top=card_top,
                  width=card_width, height=card_height,
                  theme=theme, color_key="primary")

        # Numbered badge
        _add_numbered_badge(slide, idx + 1,
                            left=card_left + Inches(0.15),
                            top=card_top + Inches(0.15),
                            fill_color=theme.colors.accent,
                            text_color="#FFFFFF",
                            font_name=theme.font_heading)

        # Title text (next to badge)
        text_left = card_left + Inches(0.15) + badge_size + Inches(0.1)
        text_width = card_width - badge_size - Inches(0.4)
        _add_text_box(slide, item_title,
                      left=text_left, top=card_top + Inches(0.15),
                      width=text_width, height=Inches(0.35),
                      font_name=theme.font_heading,
                      font_size=_DESIGN_TOKENS["body_size"],
                      font_color=theme.colors.text, bold=True)

        # Description text (below badge)
        if item_desc:
            _add_text_box(slide, item_desc,
                          left=card_left + Inches(0.2),
                          top=card_top + Inches(0.65),
                          width=card_width - Inches(0.4),
                          height=card_height - Inches(0.75),
                          font_name=theme.font_body,
                          font_size=_DESIGN_TOKENS["body_small_size"],
                          font_color=theme.colors.text_light)

    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


def _style_table_cell(cell, theme, *, is_header, is_alt):
    """Style a single table cell based on header/alternating row status."""
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(
            theme.colors.primary.lstrip("#")
        )
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.font.size = _DESIGN_TOKENS["body_small_size"]
            p.font.name = theme.font_heading
    else:
        if is_alt:
            tint = _blend_color(theme.colors.primary, theme.colors.background, 0.06)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(tint.lstrip("#"))
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                theme.colors.background.lstrip("#")
            )
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = RGBColor.from_string(
                theme.colors.text.lstrip("#")
            )
            p.font.size = _DESIGN_TOKENS["body_small_size"]
            p.font.name = theme.font_body

    # Cell margins
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


def _render_table(slide, slide_data, theme, **kwargs):
    """Table slide: styled data table with header row and alternating bands."""
    _render_kicker(slide, slide_data, theme)

    body_top, title_font, title_top = _compute_body_top(slide_data.title, has_kicker=_has_kicker(slide_data))

    _add_text_box(slide, slide_data.title or "",
                  left=TITLE_LEFT, top=title_top,
                  width=TITLE_WIDTH, height=TITLE_HEIGHT,
                  font_name=theme.font_heading, font_size=title_font,
                  font_color=theme.colors.primary, bold=True)

    table_el = _find_element(slide_data, "table_data")
    source_el = _find_element(slide_data, "source_citation")

    table_data = {}
    if table_el and isinstance(table_el.content, dict):
        table_data = table_el.content
    elif table_el and isinstance(table_el.content, str):
        try:
            import json
            table_data = json.loads(table_el.content)
        except (json.JSONDecodeError, TypeError):
            pass

    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    badge_column = table_data.get("badge_column")
    if badge_column is not None:
        try:
            badge_column = int(badge_column)
        except (ValueError, TypeError):
            badge_column = None

    if not headers:
        # Fallback: render as plain text
        _add_text_box(slide, "[Table: no data]",
                      left=BODY_LEFT, top=body_top,
                      width=BODY_WIDTH, height=Inches(2.0),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["body_small_size"],
                      font_color=theme.colors.secondary,
                      alignment=PP_ALIGN.CENTER)
        _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))
        return

    # Reserve space for source citation and takeaway
    table_bottom = SLIDE_HEIGHT - MARGIN_BOTTOM
    if source_el and source_el.content:
        table_bottom -= Inches(0.5)
    if _has_takeaway(slide_data):
        table_bottom -= _TAKEAWAY_RESERVE

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_height = min(table_bottom - body_top, Inches(4.5))

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        BODY_LEFT, body_top, BODY_WIDTH, table_height,
    )
    table = table_shape.table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        _style_table_cell(cell, theme, is_header=True, is_alt=False)

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        is_alt = row_idx % 2 == 1
        for col_idx in range(num_cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell_value = row_data[col_idx] if col_idx < len(row_data) else ""
            cell.text = str(cell_value)
            _style_table_cell(cell, theme, is_header=False, is_alt=is_alt)

    # Badge column: add tag badges overlaid on cells
    if badge_column is not None and 0 <= badge_column < num_cols:
        for row_idx, row_data in enumerate(rows):
            if badge_column < len(row_data) and row_data[badge_column]:
                badge_text = str(row_data[badge_column])
                # Clear the cell text since we'll overlay a badge
                cell = table.cell(row_idx + 1, badge_column)
                cell.text = ""
                # Approximate cell position for badge overlay
                col_width = BODY_WIDTH / num_cols
                badge_left = BODY_LEFT + badge_column * col_width + Inches(0.1)
                row_height = table_height / num_rows
                badge_top = body_top + (row_idx + 1) * row_height + Inches(0.05)
                _add_tag_badge(slide, badge_text, badge_left, badge_top,
                               fill_color=theme.colors.accent,
                               text_color="#FFFFFF",
                               font_name=theme.font_body)

    # Source citation
    if source_el and source_el.content:
        citation_top = body_top + table_height + Inches(0.1)
        _add_text_box(slide, str(source_el.content),
                      left=BODY_LEFT, top=citation_top,
                      width=BODY_WIDTH, height=Inches(0.3),
                      font_name=theme.font_body,
                      font_size=_DESIGN_TOKENS["footer_size"],
                      font_color=theme.colors.text_light,
                      alignment=PP_ALIGN.RIGHT)

    _render_takeaway(slide, slide_data, theme)
    _set_slide_background(slide, theme, (getattr(slide_data, 'metadata', None) or {}).get('slide_style') or (getattr(slide_data, 'metadata', None) or {}).get('visual_style'))


# Renderer dispatch table
_RENDERERS = {
    "title": _render_title,
    "content": _render_content,
    "two_column": _render_two_column,
    "comparison": _render_comparison,
    "timeline": _render_timeline,
    "chart": _render_chart,
    "image_text": _render_image_text,
    "image_full": _render_image_full,
    "quote": _render_quote,
    "code": _render_code,
    "team": _render_team,
    "stat": _render_stat,
    "stats": _render_stat,
    "section_divider": _render_section_divider,
    "agenda": _render_agenda,
    "table": _render_table,
}
